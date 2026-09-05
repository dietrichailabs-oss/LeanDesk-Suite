from __future__ import annotations

import tkinter as tk

import pytest

from leandesk.sheets import SHEET_CANVAS, SelectionRange, SheetGrid, SheetModel, WorkbookModel
from leandesk.core import AppSettings, RecentFiles
from leandesk.writer import WriterFrame
from leandesk.slides import DeckModel, SlideModel, SlideObject, SlidesFrame
from leandesk.draw import DrawFrame, Drawing, Shape
from leandesk.organizer import CalendarEvent, Contact, Task, _item_dict, _validated_item
from leandesk.notes import Note, NotesFrame
from leandesk.workspace import SearchResult, TemplateDefinition, TemplateStore, search_records
from leandesk.ui import apply_suite_theme


def _event(x: float, y: float, *, state: int = 0):
    return type("Event", (), {"x": x, "y": y, "state": state})()


@pytest.fixture(scope="module")
def gui_root():
    root = tk.Tk()
    root.geometry("900x500+40+40")
    yield root
    root.destroy()


def test_selection_range_normalizes_reverse_drag_and_lists_addresses() -> None:
    selection = SelectionRange("D4", "B2")
    assert selection.bounds() == (1, 1, 3, 3)
    assert selection.label() == "B2:D4"
    assert selection.addresses()[0] == "B2"
    assert selection.addresses()[-1] == "D4"
    assert len(selection.addresses()) == 9


def test_formula_translation_preserves_relative_copy_semantics() -> None:
    assert SheetGrid._translate_formula("=A1+B2", 2, 1) == "=B3+C4"
    assert SheetGrid._translate_formula("plain text", 2, 1) == "plain text"


def test_common_formula_functions_and_display_formats() -> None:
    sheet = SheetModel(cells={"A1": "-2.345", "A2": "7", "B1": "=ABS(A1)", "B2": "=ROUND(A1,1)", "B3": "=COUNTA(A1,A2)"})
    assert sheet.value("B1") == 2.345
    assert sheet.value("B2") == -2.3
    assert sheet.value("B3") == 2
    sheet.set_format("A2", number_format="currency", bold=True)
    assert sheet.display_value("A2") == "$7.00"


def test_sheet_metadata_round_trip_and_structural_operations() -> None:
    sheet = SheetModel(cells={"A1": "header", "A2": "2", "A3": "1", "B2": "two", "B3": "one"})
    sheet.set_format("A2", bold=True, align="right")
    sheet.comments["A2"] = "important"
    sheet.validations["A2"] = {"type": "number", "minimum": 0}
    sheet.freeze_panes = "B2"
    sheet.insert_rows(1)
    assert sheet.raw("A3") == "2"
    assert sheet.comments["A3"] == "important"
    sheet.delete_rows(1)
    sheet.insert_columns(1)
    assert sheet.raw("C2") == "two"
    sheet.delete_columns(1)
    sheet.sort_range("A2", "B3")
    assert sheet.raw("A2") == "1"
    assert sheet.raw("B2") == "one"
    loaded = WorkbookModel.from_dict(WorkbookModel("Book", [sheet]).to_dict()).sheets[0]
    assert loaded.cell_formats["A2"]["bold"] is True
    assert loaded.freeze_panes == "B2"


def test_grid_drag_shift_headers_all_clipboard_and_undo_redo(gui_root) -> None:
    root = gui_root
    try:
        model = SheetModel(cells={"A1": "1", "B1": "=A1+1"})
        changes: list[bool] = []
        grid = SheetGrid(root, model, lambda: changes.append(True), lambda *_args: None)
        grid.pack(fill="both", expand=True)
        root.update()

        a1 = grid._cell_bounds(0, 0)
        c3 = grid._cell_bounds(2, 2)
        grid._button_press(_event(a1[0] + 4, a1[1] + 4))
        grid._drag_motion(_event(c3[0] + 4, c3[1] + 4))
        grid._button_release(_event(c3[0] + 4, c3[1] + 4))
        assert grid.selection.label() == "A1:C3"

        grid.select_address("B2")
        grid._move_selection(2, 2, extend=True)
        assert grid.selection.label() == "B2:D4"

        grid._button_press(_event(grid._x_positions[2] + 8, 4))
        assert grid.selection.label() == "C1:C200"
        row_y = grid.HEADER_HEIGHT + grid.ROW_HEIGHT * 4 + 4
        grid._button_press(_event(4, row_y))
        assert grid.selection.label() == "A5:AZ5"
        grid.select_all()
        assert len(grid.selected_addresses()) == grid.ROWS * grid.COLS

        grid.select_range("A1", "B1")
        grid.copy_selection()
        assert root.clipboard_get() == "1\t=A1+1"
        grid.select_address("C3")
        grid.paste_selection()
        assert model.raw("C3") == "1"
        assert model.raw("D3") == "=C3+1"
        grid.undo()
        assert model.raw("C3") == ""
        assert model.raw("D3") == ""
        grid.redo()
        assert model.raw("C3") == "1"
        assert model.raw("D3") == "=C3+1"
        assert changes
    finally:
        grid.destroy()


def test_dark_theme_keeps_white_sheet_canvas_and_visible_boundaries(gui_root) -> None:
    root = gui_root
    root.withdraw()
    try:
        apply_suite_theme(root, "Dark")
        grid = SheetGrid(root, SheetModel(), lambda: None, lambda *_args: None)
        grid.pack()
        root.update()
        assert grid.canvas.cget("background") == SHEET_CANVAS
        cell = next(
            item
            for item in grid.canvas.find_withtag("grid")
            if grid.canvas.type(item) == "rectangle"
            and tuple(grid.canvas.coords(item)) == tuple(grid._cell_bounds(3, 3))
        )
        assert grid.canvas.itemcget(cell, "fill") == SHEET_CANVAS
        assert grid.canvas.itemcget(cell, "outline") != SHEET_CANVAS
        active = [
            item
            for item in grid.canvas.find_withtag("grid")
            if grid.canvas.type(item) == "rectangle"
            and float(grid.canvas.itemcget(item, "width")) == 2.0
        ]
        assert len(active) == 1
    finally:
        grid.destroy()


def test_writer_structured_table_page_parts_and_docx_round_trip(gui_root, tmp_path) -> None:
    root = gui_root
    frame = WriterFrame(root, recent=RecentFiles(tmp_path / "recent.json"), settings=AppSettings())
    frame.pack(fill="both", expand=True)
    try:
        frame.text.insert("1.0", "Quarterly report")
        object_id = frame.insert_table(2, 2, [["Item", "Amount"], ["Sales", "125"]])
        frame.set_header("Dietrich AI Labs")
        frame.set_footer("Confidential")
        frame.toggle_page_numbers()
        frame.toggle_orientation()
        document = frame.serialize()
        table = next(item for item in document.metadata["objects"] if item["id"] == object_id)
        assert table["kind"] == "table"
        assert table["data"][1] == ["Sales", "125"]
        assert document.metadata["page_numbers"] is True
        assert document.metadata["orientation"] == "landscape"

        target = tmp_path / "writer-structured.docx"
        WriterFrame._save_docx(document, target)
        loaded = WriterFrame._load_docx(target)
        assert loaded.metadata["header"] == "Dietrich AI Labs"
        assert "Confidential" in loaded.metadata["footer"]
        loaded_table = next(item for item in loaded.metadata["objects"] if item["kind"] == "table")
        assert loaded_table["data"][1] == ["Sales", "125"]
    finally:
        frame.destroy()


def test_slides_object_model_round_trip_geometry_and_powerpoint_export(tmp_path) -> None:
    slide = SlideModel(title="Results", body="Quarterly performance", layout="Two Content", transition="Fade")
    slide.objects.extend([
        SlideObject(kind="text", x=80, y=150, text="Highlights", font_size=24),
        SlideObject(kind="shape", x=500, y=150, text="Approved"),
        SlideObject(kind="table", x=80, y=300, width=360, height=120, data={"rows": 2, "cols": 2, "values": [["Item", "Value"], ["Sales", "125"]]}),
        SlideObject(kind="chart", x=500, y=280, width=360, height=180, data={"categories": ["Q1", "Q2"], "values": [10, 14]}),
    ])
    loaded = DeckModel.from_dict(DeckModel("Review", [slide]).to_dict())
    assert loaded.slides[0].layout == "Two Content"
    assert loaded.slides[0].transition == "Fade"
    assert [item.kind for item in loaded.slides[0].objects] == ["text", "shape", "table", "chart"]

    target = tmp_path / "objects.pptx"
    frame = object.__new__(SlidesFrame)
    frame.deck = loaded
    frame._save_pptx(target)
    from pptx import Presentation
    presentation = Presentation(target)
    shapes = presentation.slides[0].shapes
    assert any(shape.has_table for shape in shapes)
    assert any(shape.has_chart for shape in shapes)


def test_slides_object_move_resize_and_z_order(gui_root, tmp_path) -> None:
    frame = SlidesFrame(gui_root, recent=RecentFiles(tmp_path / "slides-recent.json"))
    frame.pack(fill="both", expand=True)
    try:
        first = frame.add_object("shape", x=100, y=100, width=200, height=100)
        second = frame.add_object("text", x=300, y=200, text="Text")
        frame.move_object(first.object_id, 50, 25)
        frame.resize_object(first.object_id, 250, 140)
        frame.reorder_object(first.object_id, 1)
        assert (first.x, first.y, first.width, first.height) == (150, 125, 250, 140)
        assert frame.current_slide().objects[-1].object_id == first.object_id
        assert second.object_id != first.object_id
    finally:
        frame.destroy()


def test_draw_multiselect_group_resize_rotate_zorder_clipboard_and_history(gui_root, tmp_path) -> None:
    frame = DrawFrame(gui_root, recent=RecentFiles(tmp_path / "draw-recent.json"))
    frame.drawing = Drawing(shapes=[Shape("a", "rectangle", 0, 0, 100, 50), Shape("b", "ellipse", 120, 0, 200, 80)])
    frame.mark_dirty = lambda: None
    try:
        frame.select_shapes(["a", "b"]); frame.group_selected()
        assert frame.drawing.shapes[0].group_id == frame.drawing.shapes[1].group_id
        frame.resize_selected(2, 2); frame.rotate_selected(90); frame.send_to_back()
        assert frame.shape_by_id("a").x2 == 200
        assert frame.shape_by_id("b").rotation == 90
        frame.copy_selected(); frame.paste_shapes()
        assert len(frame.drawing.shapes) == 4
        frame.undo(); assert len(frame.drawing.shapes) == 2
        frame.redo(); assert len(frame.drawing.shapes) == 4
        loaded = Drawing.from_dict(frame.drawing.to_dict())
        assert loaded.shapes[0].group_id
    finally:
        frame.destroy()


def test_organizer_rich_fields_round_trip_without_losing_extensions() -> None:
    task = Task("12345678-1234-5678-1234-567812345678", "Review", recurrence="Weekly", reminder="09:00", categories="QA")
    event = CalendarEvent("12345678-1234-5678-1234-567812345679", "2026-09-03", "Launch", location="HQ", recurrence="Yearly", all_day="true")
    contact = Contact("12345678-1234-5678-1234-567812345680", "Mark", job_title="Founder", website="https://example.com", categories="Partner")
    for model, item, name in ((Task, task, "tasks"), (CalendarEvent, event, "events"), (Contact, contact, "contacts")):
        payload = _item_dict(item); payload["future_field"] = "preserved"
        loaded = _validated_item(model, payload, collection=name)
        assert _item_dict(loaded)["future_field"] == "preserved"
    assert _validated_item(Task, _item_dict(task), collection="tasks").recurrence == "Weekly"
    assert _validated_item(CalendarEvent, _item_dict(event), collection="events").location == "HQ"
    assert _validated_item(Contact, _item_dict(contact), collection="contacts").job_title == "Founder"


def test_notes_hierarchy_and_structured_markdown_round_trip() -> None:
    parent = Note.new("Project")
    child = Note.new("Meeting")
    child.parent_id = parent.note_id
    child.body = "- [ ] Follow up\n\n" + NotesFrame.markdown_table(2, 3)
    loaded = Note.from_dict(child.to_dict())
    assert loaded.parent_id == parent.note_id
    assert "- [ ] Follow up" in loaded.body
    assert "| Column 1 | Column 2 | Column 3 |" in loaded.body


def test_template_store_and_global_search_are_bounded_and_round_trip(tmp_path) -> None:
    store = TemplateStore(tmp_path / "templates")
    custom = TemplateDefinition("QA Checklist", "Notes", {"title": "QA", "body": "- [ ] Verify"})
    store.save(custom)
    assert any(row.name == "QA Checklist" for row in store.list())
    rows = [SearchResult("Notes", "QA Checklist", "Verify installer signatures", "1"), SearchResult("Tasks", "Ship", "Publish release", "2")]
    assert search_records("qa signatures", rows) == [rows[0]]
    assert search_records("", rows) == []

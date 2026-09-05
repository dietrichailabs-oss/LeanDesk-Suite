from __future__ import annotations

import io
import json
import os
from pathlib import Path
import shutil
import stat
import subprocess
from types import SimpleNamespace
from unittest.mock import patch
import zipfile
from xml.etree import ElementTree as ET

import pytest

from leandesk.backup_service import (
    BackupRestoreStateError,
    create_backup,
    recover_abandoned_restore_state,
    restore_backup,
)
from leandesk.document_formats import LeanDocument, plain_to_rtf, read_text_document, write_text_document
from leandesk.ooxml_preflight import (
    OOXMLLimits,
    OOXMLPreflightCancelled,
    OOXMLPreflightError,
    OOXMLPreflightTimeout,
    prepare_ooxml,
)
from leandesk.rtf_codec import RTFFormatError, rtf_to_plain
from leandesk.updates.update_manifest import ManifestError, parse_manifest
from leandesk.updates.version_compare import Version, VersionError, compare_versions

ROOT = Path(__file__).resolve().parents[1]
FIXTURES = Path(__file__).resolve().parent / "fixtures" / "correction_3"
OFFICE_DUPLICATES = FIXTURES / "OFFICE_DUPLICATE_PARTS"
RTF_FIXTURES = FIXTURES / "RTF"


def _manifest_bytes(version: str) -> bytes:
    return json.dumps(
        {
            "product": "leandesk-suite",
            "latest_version": version,
            "release_name": f"LeanDesk Suite {version}",
            "published_at": "2026-08-24T00:00:00Z",
            "release_url": "https://www.dietrichailabs.com/apps/leandesk/",
            "download_url": "https://www.dietrichailabs.com/downloads/",
            "sha256": "",
            "message": "A new version is available.",
        }
    ).encode("utf-8")


def _make_profile_backup(tmp_path: Path) -> tuple[Path, Path, Path]:
    source = tmp_path / "source_profile"
    live = tmp_path / "live_profile"
    archive = tmp_path / "profile.ldbackup"
    source.mkdir()
    live.mkdir()
    (source / "state.txt").write_text("NEW_PROFILE", encoding="utf-8")
    (live / "state.txt").write_text("OLD_PROFILE", encoding="utf-8")
    create_backup(archive, data_root=source)
    return source, live, archive


def _make_valid_office(tmp_path: Path, kind: str) -> Path:
    path = tmp_path / f"valid.{kind}"
    if kind == "docx":
        from docx import Document

        document = Document()
        document.add_paragraph("VALID_DOCX")
        document.save(path)
    elif kind == "xlsx":
        from openpyxl import Workbook

        workbook = Workbook()
        workbook.active["A1"] = "VALID_XLSX"
        workbook.save(path)
    elif kind == "pptx":
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "VALID_PPTX"
        presentation.save(path)
    else:
        raise ValueError(kind)
    return path


def _rewrite_zip(
    source: Path,
    destination: Path,
    *,
    replacements: dict[str, bytes] | None = None,
    additions: list[tuple[str | zipfile.ZipInfo, bytes]] | None = None,
) -> Path:
    replacements = replacements or {}
    additions = additions or []
    with zipfile.ZipFile(source, "r") as incoming, zipfile.ZipFile(
        destination, "w", compression=zipfile.ZIP_DEFLATED
    ) as outgoing:
        for info in incoming.infolist():
            data = incoming.read(info)
            outgoing.writestr(info, replacements.get(info.filename, data))
        for name, data in additions:
            if isinstance(name, str) and "\\" in name:
                # ZipInfo(string) normalizes os.sep on Windows. Assigning filename
                # after construction preserves the raw ambiguous member for this
                # cross-platform security regression.
                raw = zipfile.ZipInfo("placeholder")
                raw.filename = name
                raw.compress_type = zipfile.ZIP_DEFLATED
                outgoing.writestr(raw, data)
            else:
                outgoing.writestr(name, data)
    return destination


# ---------------------------------------------------------------------------
# LD-QA-C2-001: one consistent version ordering relation
# ---------------------------------------------------------------------------


def test_version_case_variants_are_antisymmetric() -> None:
    forward = compare_versions("1.0.0-ALPHA", "1.0.0-alpha")
    reverse = compare_versions("1.0.0-alpha", "1.0.0-ALPHA")
    assert forward == -reverse
    assert forward != 0


@pytest.mark.parametrize("value", ["1.0.0-00", "1.0.0-01", "1.0.0-alpha.01", "2.3.4-000.beta"])
def test_version_rejects_leading_zero_numeric_prerelease(value: str) -> None:
    with pytest.raises(VersionError):
        Version.parse(value)


def test_version_generated_matrix_is_antisymmetric_and_equality_consistent() -> None:
    versions = [
        "0.8.0-alpha",
        "0.8.0-alpha.1",
        "0.8.0-ALPHA",
        "0.8.0-beta.2",
        "0.8.0-rc.1",
        "0.8.0",
        "0.8.0.1",
        "0.9.0",
        "0.10.0",
        "1.0.0+build.1",
        "1.0.0+build.2",
    ]
    for left in versions:
        for right in versions:
            ab = compare_versions(left, right)
            ba = compare_versions(right, left)
            assert ab in {-1, 0, 1}
            assert ab == -ba
            assert (ab == 0) is (Version.parse(left) == Version.parse(right))


def test_version_semver_precedence_chain_is_transitive() -> None:
    ordered = [
        "1.0.0-alpha",
        "1.0.0-alpha.1",
        "1.0.0-alpha.beta",
        "1.0.0-beta",
        "1.0.0-beta.2",
        "1.0.0-beta.11",
        "1.0.0-rc.1",
        "1.0.0",
        "1.0.1",
        "1.1.0",
        "2.0.0",
    ]
    for first_index, first in enumerate(ordered):
        for second_index in range(first_index + 1, len(ordered)):
            second = ordered[second_index]
            assert compare_versions(first, second) < 0
            for third in ordered[second_index + 1 :]:
                assert compare_versions(first, third) < 0


def test_version_build_metadata_is_ignored_for_precedence_and_equality() -> None:
    assert compare_versions("1.2.3+build.1", "1.2.3+build.999") == 0
    assert Version.parse("1.2.3+build.1") == Version.parse("1.2.3+build.999")


def test_update_manifest_fails_closed_on_ambiguous_prerelease() -> None:
    with pytest.raises(ManifestError):
        parse_manifest(_manifest_bytes("1.0.0-01"))


# ---------------------------------------------------------------------------
# LD-QA-C2-002: transactionally truthful restore and restart recovery
# ---------------------------------------------------------------------------


def test_restore_cleanup_failure_after_commit_returns_truthful_success(tmp_path: Path) -> None:
    _source, live, archive = _make_profile_backup(tmp_path)
    real_rmtree = shutil.rmtree
    failed = False

    def fail_first_rollback_cleanup(path, *args, **kwargs):
        nonlocal failed
        candidate = Path(path)
        if ".live_profile.rollback-" in candidate.name and not failed:
            failed = True
            raise OSError("simulated rollback cleanup failure")
        return real_rmtree(path, *args, **kwargs)

    with patch("leandesk.backup_service.shutil.rmtree", side_effect=fail_first_rollback_cleanup):
        result = restore_backup(archive, data_root=live)

    assert result["committed"] is True
    assert (live / "state.txt").read_text(encoding="utf-8") == "NEW_PROFILE"
    assert result["rollback_retained"] is True
    assert Path(result["rollback_path"]).is_dir()
    assert "restored profile is live" in result["cleanup_warning"]


def test_restore_first_rename_failure_keeps_old_profile_active(tmp_path: Path) -> None:
    _source, live, archive = _make_profile_backup(tmp_path)
    real_replace = os.replace

    def fail_live_move(source, destination):
        if Path(source) == live:
            raise OSError("simulated first rename failure")
        return real_replace(source, destination)

    with patch("leandesk.backup_service.os.replace", side_effect=fail_live_move):
        with pytest.raises(BackupRestoreStateError) as caught:
            restore_backup(archive, data_root=live)
    assert caught.value.profile_state == "profile_unchanged"
    assert (live / "state.txt").read_text(encoding="utf-8") == "OLD_PROFILE"
    assert not list(tmp_path.glob(".live_profile.rollback-*"))


def test_restore_second_rename_failure_rolls_old_profile_back(tmp_path: Path) -> None:
    _source, live, archive = _make_profile_backup(tmp_path)
    real_replace = os.replace
    calls = 0

    def fail_second_replace(source, destination):
        nonlocal calls
        calls += 1
        if calls == 2:
            raise OSError("simulated staging-to-live failure")
        return real_replace(source, destination)

    with patch("leandesk.backup_service.os.replace", side_effect=fail_second_replace):
        with pytest.raises(BackupRestoreStateError) as caught:
            restore_backup(archive, data_root=live)
    assert caught.value.profile_state == "previous_profile_active"
    assert (live / "state.txt").read_text(encoding="utf-8") == "OLD_PROFILE"
    assert not list(tmp_path.glob(".live_profile.rollback-*"))


def test_restore_precommit_staging_fsync_failure_keeps_old_profile(tmp_path: Path) -> None:
    _source, live, archive = _make_profile_backup(tmp_path)
    with patch("leandesk.backup_service._sync_staging_tree", side_effect=OSError("fsync")):
        with pytest.raises(BackupRestoreStateError) as caught:
            restore_backup(archive, data_root=live)
    assert caught.value.profile_state == "profile_unchanged"
    assert (live / "state.txt").read_text(encoding="utf-8") == "OLD_PROFILE"


def test_restore_sync_failure_between_renames_restores_old_profile(tmp_path: Path) -> None:
    _source, live, archive = _make_profile_backup(tmp_path)
    calls = 0

    def sync_once_fails(_path):
        nonlocal calls
        calls += 1
        if calls == 1:
            raise OSError("parent fsync after first rename")

    with (
        patch("leandesk.backup_service._sync_staging_tree", return_value=None),
        patch("leandesk.backup_service._sync_directory", side_effect=sync_once_fails),
    ):
        with pytest.raises(BackupRestoreStateError) as caught:
            restore_backup(archive, data_root=live)
    assert caught.value.profile_state == "previous_profile_active"
    assert (live / "state.txt").read_text(encoding="utf-8") == "OLD_PROFILE"


def test_restore_postcommit_fsync_failure_is_success_with_warning(tmp_path: Path) -> None:
    _source, live, archive = _make_profile_backup(tmp_path)
    calls = 0

    def fail_second_sync(_path):
        nonlocal calls
        calls += 1
        if calls == 2:
            raise OSError("postcommit parent fsync")

    with (
        patch("leandesk.backup_service._sync_staging_tree", return_value=None),
        patch("leandesk.backup_service._sync_directory", side_effect=fail_second_sync),
    ):
        result = restore_backup(archive, data_root=live)
    assert result["committed"] is True
    assert (live / "state.txt").read_text(encoding="utf-8") == "NEW_PROFILE"
    assert "durability flush failed" in result["cleanup_warning"]


def test_restore_postcommit_staging_cleanup_failure_is_success(tmp_path: Path) -> None:
    _source, live, archive = _make_profile_backup(tmp_path)
    real_rmtree = shutil.rmtree

    def fail_staging_cleanup(path, *args, **kwargs):
        if ".live_profile.restore-staging-" in Path(path).name:
            raise OSError("staging cleanup")
        return real_rmtree(path, *args, **kwargs)

    with patch("leandesk.backup_service.shutil.rmtree", side_effect=fail_staging_cleanup):
        result = restore_backup(archive, data_root=live)
    assert result["committed"] is True
    assert (live / "state.txt").read_text(encoding="utf-8") == "NEW_PROFILE"
    assert "restore-staging directory" in result["cleanup_warning"]


def test_restore_rollback_reactivation_failure_retains_old_profile_copy(tmp_path: Path) -> None:
    _source, live, archive = _make_profile_backup(tmp_path)
    real_replace = os.replace
    calls = 0

    def fail_new_commit_and_rollback_reactivation(source, destination):
        nonlocal calls
        calls += 1
        if calls in {2, 3}:
            raise OSError("simulated rename boundary failure")
        return real_replace(source, destination)

    with patch(
        "leandesk.backup_service.os.replace",
        side_effect=fail_new_commit_and_rollback_reactivation,
    ):
        with pytest.raises(BackupRestoreStateError) as caught:
            restore_backup(archive, data_root=live)
    assert caught.value.profile_state == "previous_profile_retained"
    assert not live.exists()
    rollbacks = list(tmp_path.glob(".live_profile.rollback-*"))
    assert len(rollbacks) == 1
    assert (rollbacks[0] / "state.txt").read_text(encoding="utf-8") == "OLD_PROFILE"


def test_restart_recovery_reactivates_sole_rollback_and_cleans_staging(tmp_path: Path) -> None:
    live = tmp_path / "profile"
    rollback = tmp_path / ".profile.rollback-deadbeef"
    staging = tmp_path / ".profile.restore-staging-deadbeef"
    rollback.mkdir()
    staging.mkdir()
    (rollback / "state.txt").write_text("OLD_PROFILE", encoding="utf-8")
    (staging / "partial.txt").write_text("UNCOMMITTED", encoding="utf-8")

    result = recover_abandoned_restore_state(data_root=live)
    assert result["recovered_previous_profile"] is True
    assert (live / "state.txt").read_text(encoding="utf-8") == "OLD_PROFILE"
    assert not rollback.exists()
    assert not staging.exists()


def test_restart_recovery_retains_rollback_when_live_profile_exists(tmp_path: Path) -> None:
    live = tmp_path / "profile"
    rollback = tmp_path / ".profile.rollback-deadbeef"
    live.mkdir()
    rollback.mkdir()
    (live / "state.txt").write_text("NEW_PROFILE", encoding="utf-8")
    (rollback / "state.txt").write_text("OLD_PROFILE", encoding="utf-8")

    result = recover_abandoned_restore_state(data_root=live)
    assert result["recovered_previous_profile"] is False
    assert result["retained_rollback_paths"] == [str(rollback)]
    assert rollback.is_dir()


def test_restart_recovery_refuses_to_guess_between_multiple_rollbacks(tmp_path: Path) -> None:
    for token in ("one", "two"):
        rollback = tmp_path / f".profile.rollback-{token}"
        rollback.mkdir()
        (rollback / "state.txt").write_text(token, encoding="utf-8")
    with pytest.raises(BackupRestoreStateError) as caught:
        recover_abandoned_restore_state(data_root=tmp_path / "profile")
    assert caught.value.profile_state == "previous_profile_retained"
    assert len(list(tmp_path.glob(".profile.rollback-*"))) == 2


def test_restore_ui_message_matches_retained_rollback_state() -> None:
    from leandesk.app import LeanDeskApp

    shown: list[str] = []
    dummy = SimpleNamespace(destroy=lambda: None)
    rollback_path = Path("/safe/rollback")
    error = BackupRestoreStateError(
        "simulated",
        profile_state="previous_profile_retained",
        rollback_path=rollback_path,
    )
    with (
        patch("leandesk.app.filedialog.askopenfilename", return_value="backup.ldbackup"),
        patch("leandesk.app.messagebox.askyesno", return_value=True),
        patch("leandesk.app.restore_backup", side_effect=error),
        patch("leandesk.app.messagebox.showerror", side_effect=lambda _title, text, **_kw: shown.append(text)),
    ):
        LeanDeskApp.restore_profile_backup(dummy)
    assert shown
    assert "preserved in a recovery directory" in shown[0]
    assert str(rollback_path) in shown[0]
    assert "remains in place" not in shown[0]


# ---------------------------------------------------------------------------
# LD-QA-C2-003: shared bounded Office archive preflight
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("kind", ["docx", "xlsx", "pptx"])
def test_supplied_conflicting_duplicate_office_packages_are_rejected(kind: str) -> None:
    with pytest.raises(OOXMLPreflightError, match="exact duplicate"):
        prepare_ooxml(OFFICE_DUPLICATES / f"conflicting-duplicate.{kind}", kind)


@pytest.mark.parametrize(
    ("kind", "patch_target", "loader_path"),
    [
        ("docx", "docx.Document", "writer"),
        ("xlsx", "openpyxl.load_workbook", "sheets"),
        ("pptx", "pptx.Presentation", "slides"),
    ],
)
def test_duplicate_package_is_rejected_before_third_party_parser(
    kind: str, patch_target: str, loader_path: str
) -> None:
    if loader_path == "writer":
        from leandesk.writer import WriterFrame

        loader = WriterFrame._load_docx
    elif loader_path == "sheets":
        from leandesk.sheets import SheetsFrame

        loader = SheetsFrame._load_xlsx
    else:
        from leandesk.slides import SlidesFrame

        loader = SlidesFrame._load_pptx

    with patch(patch_target) as parser:
        with pytest.raises(OOXMLPreflightError):
            loader(OFFICE_DUPLICATES / f"conflicting-duplicate.{kind}")
        parser.assert_not_called()


@pytest.mark.parametrize("kind", ["docx", "xlsx", "pptx"])
def test_valid_office_packages_pass_preflight_and_actual_loader(tmp_path: Path, kind: str) -> None:
    path = _make_valid_office(tmp_path, kind)
    report = prepare_ooxml(path, kind).report
    assert report.entries > 0
    if kind == "docx":
        from leandesk.writer import WriterFrame

        assert "VALID_DOCX" in WriterFrame._load_docx(path).text
    elif kind == "xlsx":
        from leandesk.sheets import SheetsFrame

        assert SheetsFrame._load_xlsx(path).sheets[0].cells["A1"] == "VALID_XLSX"
    else:
        from leandesk.slides import SlidesFrame

        assert SlidesFrame._load_pptx(path).slides[0].title == "VALID_PPTX"


@pytest.mark.parametrize(
    "unsafe_name",
    ["../evil.xml", "/absolute.xml", "C:/drive.xml", "word\\ambiguous.xml", "folder/CON.xml"],
)
def test_ooxml_rejects_unsafe_member_paths(tmp_path: Path, unsafe_name: str) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    malicious = _rewrite_zip(valid, tmp_path / "unsafe.docx", additions=[(unsafe_name, b"<x/>")])
    with pytest.raises(OOXMLPreflightError):
        prepare_ooxml(malicious, "docx")


def test_ooxml_rejects_case_colliding_members(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    with zipfile.ZipFile(valid, "r") as archive:
        payload = archive.read("word/document.xml")
    malicious = _rewrite_zip(
        valid,
        tmp_path / "case-collision.docx",
        additions=[("WORD/document.xml", payload)],
    )
    with pytest.raises(OOXMLPreflightError, match="case/Unicode-colliding"):
        prepare_ooxml(malicious, "docx")


def test_ooxml_rejects_link_or_special_member(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    link = zipfile.ZipInfo("word/media/link")
    link.create_system = 3
    link.external_attr = (stat.S_IFLNK | 0o777) << 16
    malicious = _rewrite_zip(valid, tmp_path / "linked.docx", additions=[(link, b"target")])
    with pytest.raises(OOXMLPreflightError, match="linked or special"):
        prepare_ooxml(malicious, "docx")


def test_ooxml_rejects_encrypted_member_flag(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    data = bytearray(valid.read_bytes())
    local = data.find(b"PK\x03\x04")
    central = data.find(b"PK\x01\x02")
    assert local >= 0 and central >= 0
    local_flags = int.from_bytes(data[local + 6 : local + 8], "little") | 0x1
    central_flags = int.from_bytes(data[central + 8 : central + 10], "little") | 0x1
    data[local + 6 : local + 8] = local_flags.to_bytes(2, "little")
    data[central + 8 : central + 10] = central_flags.to_bytes(2, "little")
    malicious = tmp_path / "encrypted.docx"
    malicious.write_bytes(data)
    with pytest.raises(OOXMLPreflightError, match="encrypted member"):
        prepare_ooxml(malicious, "docx")


def test_ooxml_rejects_excessive_entries_with_configurable_budget(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    with pytest.raises(OOXMLPreflightError, match="too many members"):
        prepare_ooxml(valid, "docx", limits=OOXMLLimits(max_entries=1))


def test_ooxml_rejects_oversized_xml_with_configurable_budget(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    with pytest.raises(OOXMLPreflightError, match="XML size limit"):
        prepare_ooxml(valid, "docx", limits=OOXMLLimits(max_xml_bytes=100))


def test_ooxml_rejects_oversized_media_with_configurable_budget(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    malicious = _rewrite_zip(
        valid,
        tmp_path / "media.docx",
        additions=[("word/media/large.bin", b"M" * 512)],
    )
    with pytest.raises(OOXMLPreflightError, match="media size limit"):
        prepare_ooxml(malicious, "docx", limits=OOXMLLimits(max_media_bytes=64))


def test_ooxml_rejects_pathological_compression_ratio(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    malicious = _rewrite_zip(
        valid,
        tmp_path / "ratio.docx",
        additions=[("word/media/zeros.bin", b"\x00" * 100_000)],
    )
    with pytest.raises(OOXMLPreflightError, match="compression-ratio"):
        prepare_ooxml(malicious, "docx", limits=OOXMLLimits(max_compression_ratio=2.0))


def test_ooxml_rejects_total_expansion_over_budget(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "xlsx")
    with pytest.raises(OOXMLPreflightError, match="total expansion"):
        prepare_ooxml(valid, "xlsx", limits=OOXMLLimits(max_total_expanded_bytes=1024))


def test_ooxml_rejects_dtd_or_entity_xml(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    with zipfile.ZipFile(valid, "r") as archive:
        content_types = archive.read("[Content_Types].xml")
    malicious = _rewrite_zip(
        valid,
        tmp_path / "dtd.docx",
        replacements={"[Content_Types].xml": b'<!DOCTYPE Types [<!ENTITY x "boom">]>' + content_types},
    )
    with pytest.raises(OOXMLPreflightError, match="forbidden DTD"):
        prepare_ooxml(malicious, "docx")


def test_ooxml_rejects_malformed_relationship_xml(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    malicious = _rewrite_zip(
        valid,
        tmp_path / "malformed-rels.docx",
        replacements={"_rels/.rels": b"<Relationships><broken"},
    )
    with pytest.raises(OOXMLPreflightError, match="malformed"):
        prepare_ooxml(malicious, "docx")


def test_ooxml_rejects_missing_internal_relationship_target(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    with zipfile.ZipFile(valid, "r") as archive:
        rels = archive.read("_rels/.rels").replace(b"word/document.xml", b"word/missing.xml")
    malicious = _rewrite_zip(
        valid,
        tmp_path / "missing-target.docx",
        replacements={"_rels/.rels": rels},
    )
    with pytest.raises(OOXMLPreflightError, match="target is missing"):
        prepare_ooxml(malicious, "docx")


def test_ooxml_rejects_unsafe_external_relationship(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    with zipfile.ZipFile(valid, "r") as archive:
        rels = archive.read("word/_rels/document.xml.rels")
    addition = (
        b'<Relationship Id="rId999" '
        b'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" '
        b'Target="file:///etc/passwd" TargetMode="External"/>'
    )
    rels = rels.replace(b"</Relationships>", addition + b"</Relationships>")
    malicious = _rewrite_zip(
        valid,
        tmp_path / "external.docx",
        replacements={"word/_rels/document.xml.rels": rels},
    )
    with pytest.raises(OOXMLPreflightError, match="Unsupported external relationship"):
        prepare_ooxml(malicious, "docx")


def test_ooxml_cancellation_stops_before_parser(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    with pytest.raises(OOXMLPreflightCancelled):
        prepare_ooxml(valid, "docx", cancel_check=lambda: True)


def test_ooxml_runtime_budget_timeout_is_controlled(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    ticks = iter((0.0, 1.0))
    with pytest.raises(OOXMLPreflightTimeout):
        prepare_ooxml(
            valid,
            "docx",
            limits=OOXMLLimits(timeout_seconds=0.05),
            clock=lambda: next(ticks, 1.0),
        )


def test_ooxml_rejects_symlink_source(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")
    alias = tmp_path / "alias.docx"
    try:
        alias.symlink_to(valid)
    except (OSError, NotImplementedError):
        pytest.skip("symlink creation unavailable")
    with pytest.raises(OOXMLPreflightError, match="link or reparse"):
        prepare_ooxml(alias, "docx")


def test_ooxml_seekable_stream_position_is_preserved(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "xlsx")
    stream = io.BytesIO(valid.read_bytes())
    stream.seek(7)
    prepared = prepare_ooxml(stream, "xlsx")
    assert stream.tell() == 7
    assert prepared.payload == valid.read_bytes()


# ---------------------------------------------------------------------------
# LD-QA-C2-004: Unicode-safe, self-round-trippable RTF
# ---------------------------------------------------------------------------


def test_supplied_standard_unicode_rtf_loads_exact_text() -> None:
    loaded = read_text_document(RTF_FIXTURES / "standard_unicode.rtf")
    assert loaded.text == "LeanDesk Ω é 😀 中文"


@pytest.mark.parametrize(
    "text",
    [
        "ASCII only",
        "Latin-1: café déjà vu £ ¥",
        "Greek: Ω Καλημέρα",
        "CJK: 中文 日本語 한국어",
        "Combining: e\u0301 A\u030a",
        "Emoji: 😀 🚀 🧪",
        "Line one\nLine two\tTabbed",
        "",
    ],
)
def test_rtf_own_output_round_trips_unicode(text: str) -> None:
    rendered = plain_to_rtf(text)
    assert rendered.isascii()
    assert rtf_to_plain(rendered.encode("ascii")) == text


def test_rtf_skips_font_color_and_info_destinations() -> None:
    payload = (
        r"{\rtf1\ansi"
        r"{\fonttbl{\f0 Segoe UI;}}"
        r"{\colortbl;\red255\green0\blue0;}"
        r"{\info{\title Hidden title}{\author Hidden author}}"
        r"Visible text}"
    ).encode("ascii")
    assert rtf_to_plain(payload) == "Visible text"


def test_rtf_decodes_ansi_hex_escape_using_declared_codepage() -> None:
    payload = rb"{\rtf1\ansi\ansicpg1252 Caf\'e9}"
    assert rtf_to_plain(payload) == "Café"


def test_rtf_honors_uc_fallback_count_and_surrogate_pairs() -> None:
    payload = rb"{\rtf1\ansi\uc2 \u937?? \u-10179??\u-8704??}"
    assert rtf_to_plain(payload) == "Ω 😀"


@pytest.mark.parametrize(
    "payload",
    [
        b"not an rtf file",
        rb"{\rtf1 unclosed",
        rb"{\rtf1 \'zz}",
        rb"{\rtf1 \u999999?}",
        rb"{\rtf1 } extra }",
    ],
)
def test_rtf_malformed_input_fails_in_controlled_way(payload: bytes) -> None:
    with pytest.raises(RTFFormatError):
        rtf_to_plain(payload)


def test_rtf_oversized_input_is_rejected_without_parsing(monkeypatch: pytest.MonkeyPatch) -> None:
    import leandesk.rtf_codec as codec

    monkeypatch.setattr(codec, "MAX_RTF_BYTES", 32)
    with pytest.raises(RTFFormatError, match="safety limit"):
        codec.rtf_to_plain(b"{\\rtf1 " + b"x" * 64 + b"}")


def test_rtf_file_write_is_ascii_and_reopens_without_metadata(tmp_path: Path) -> None:
    text = "LeanDesk Ω é 😀 中文"
    target = tmp_path / "unicode.rtf"
    write_text_document(LeanDocument(title="Unicode", text=text), target)
    raw = target.read_bytes()
    assert raw.isascii()
    assert b"Segoe UI" in raw
    assert b"\xce\xa9" not in raw
    assert read_text_document(target).text == text


def test_leandesk_rtf_is_readable_by_libreoffice(tmp_path: Path) -> None:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        pytest.skip("LibreOffice is not installed")
    source_text = "LeanDesk Ω é 😀 中文 e\u0301"
    source = tmp_path / "unicode.rtf"
    outdir = tmp_path / "out"
    profile = tmp_path / "lo-profile"
    outdir.mkdir()
    profile.mkdir()
    write_text_document(LeanDocument(title="Unicode", text=source_text), source)
    command = [
        soffice,
        "--headless",
        "--nologo",
        "--nodefault",
        "--nolockcheck",
        "--norestore",
        f"-env:UserInstallation={profile.resolve().as_uri()}",
        "--convert-to",
        "txt:Text",
        "--outdir",
        str(outdir),
        str(source),
    ]
    result = subprocess.run(command, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=30, check=False)
    converted = outdir / "unicode.txt"
    assert result.returncode == 0, (result.stdout + result.stderr).decode("utf-8", errors="replace")
    assert converted.is_file()
    readback = converted.read_text(encoding="utf-8-sig", errors="strict")
    for expected in ("LeanDesk", "Ω", "é", "😀", "中文", "e\u0301"):
        assert expected in readback


@pytest.mark.skipif(not os.environ.get("DISPLAY") and os.name != "nt", reason="GUI test requires an active display/Xvfb")
def test_writer_actual_gui_rtf_open_and_save_as_round_trip(tmp_path: Path) -> None:
    import tkinter as tk

    from leandesk.core import AppSettings, RecentFiles, RecoveryStore
    from leandesk.writer import WriterFrame

    source = tmp_path / "standard_unicode.rtf"
    source.write_bytes((RTF_FIXTURES / "standard_unicode.rtf").read_bytes())
    destination = tmp_path / "saved_unicode.rtf"
    root = tk.Tk()
    root.withdraw()
    try:
        frame = WriterFrame(root, recent=RecentFiles(tmp_path / "recent.json"), settings=AppSettings())
        frame.recovery = RecoveryStore(tmp_path / "Recovery")
        with patch("leandesk.writer.messagebox.showerror") as showerror:
            assert frame.open_document(source) is True
            assert frame.text.get("1.0", "end-1c") == "LeanDesk Ω é 😀 中文"
            with patch("leandesk.writer.filedialog.asksaveasfilename", return_value=str(destination)):
                assert frame.save_as() is True
            showerror.assert_not_called()
        assert read_text_document(destination).text == "LeanDesk Ω é 😀 中文"
    finally:
        root.destroy()


def test_correction_3_modules_are_in_authoritative_test_gate() -> None:
    source = (ROOT / "tools" / "run_authoritative_tests.py").read_text(encoding="utf-8")
    assert '"tests/test_correction_3_qa.py"' in source

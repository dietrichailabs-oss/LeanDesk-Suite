from __future__ import annotations

"""Mandatory Correction 5 acceptance and adversarial regression suite.

These tests convert the exact Independent-QA Correction 4 findings into canonical
acceptance gates while varying timing, archive identity, OOXML content-type mapping,
DBCS lexical boundaries, and filesystem containment.
"""

import hashlib
import os
from pathlib import Path, PurePosixPath
import re
import shutil
import stat
import subprocess
from unittest.mock import patch
import xml.etree.ElementTree as ET
import zipfile

import pytest

from leandesk import backup_service
from leandesk.backup_integrity import BackupIntegrityError, sha256_file, verify_backup_artifact
from leandesk.backup_service import BackupRestoreStateError, create_backup, restore_backup
from leandesk.document_formats import read_text_document, write_text_document
from leandesk.ooxml_preflight import OOXMLPreflightError, prepare_ooxml
from leandesk.rtf_codec import RTFFormatError, rtf_to_plain

ROOT = Path(__file__).resolve().parents[1]
_CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------


def _profile(path: Path, text: str) -> Path:
    path.mkdir(parents=True)
    (path / "state.txt").write_text(text, encoding="utf-8")
    return path


def _backup(path: Path, profile: Path) -> Path:
    create_backup(path, data_root=profile)
    assert path.is_file()
    return path


def _restore_fixture(tmp_path: Path) -> tuple[Path, Path, str]:
    profile = _profile(tmp_path / "source_profile", "BACKUP_A_VERIFIED")
    archive = _backup(tmp_path / "source.ldbackup", profile)
    live = _profile(tmp_path / "live", "LIVE_OLD")
    return archive, live, sha256_file(archive)


def _assert_live(live: Path, expected: str) -> None:
    assert (live / "state.txt").read_text(encoding="utf-8") == expected


def _old_destination_and_new_profile(tmp_path: Path) -> tuple[Path, Path, str]:
    old_profile = _profile(tmp_path / "old_profile", "OLD_VALID_BACKUP")
    destination = _backup(tmp_path / "selected.ldbackup", old_profile)
    old_sha = sha256_file(destination)
    new_profile = _profile(tmp_path / "new_profile", "NEW_PROFILE")
    return destination, new_profile, old_sha


def _assert_old_destination(destination: Path, old_sha: str) -> None:
    assert destination.is_file()
    assert sha256_file(destination) == old_sha
    assert verify_backup_artifact(destination, require_manifest=True)["valid"] is True


def _make_valid_office(tmp_path: Path, kind: str) -> Path:
    path = tmp_path / f"valid.{kind}"
    if kind == "docx":
        from docx import Document

        document = Document()
        document.add_paragraph("VISIBLE_DOCX")
        document.save(path)
    elif kind == "xlsx":
        from openpyxl import Workbook

        workbook = Workbook()
        workbook.active["A1"] = "VISIBLE_XLSX"
        workbook.save(path)
    elif kind == "pptx":
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "VISIBLE_PPTX"
        presentation.save(path)
    else:  # pragma: no cover - helper contract
        raise ValueError(kind)
    return path


_OFFICE_PARTS = {
    "docx": {
        "part": "word/styles.xml",
        "root": "w:styles",
        "relationships": [("word/_rels/document.xml.rels", "styles.xml")],
    },
    "xlsx": {
        "part": "xl/styles.xml",
        "root": "styleSheet",
        "relationships": [("xl/_rels/workbook.xml.rels", "styles.xml")],
    },
    "pptx": {
        "part": "ppt/theme/theme1.xml",
        "root": "a:theme",
        "relationships": [
            ("ppt/_rels/presentation.xml.rels", "theme/theme1.xml"),
            ("ppt/slideMasters/_rels/slideMaster1.xml.rels", "../theme/theme1.xml"),
        ],
    },
}


def _insert_dtd(payload: bytes, root_name: str) -> bytes:
    declaration = f'<!DOCTYPE {root_name} [<!ENTITY leandesk "blocked">]>'.encode("ascii")
    if payload.lstrip().startswith(b"<?xml") and b"?>" in payload:
        offset = payload.index(b"?>") + 2
        return payload[:offset] + declaration + payload[offset:]
    return declaration + payload


def _content_types_for_renamed_part(
    payload: bytes,
    *,
    old_part: str,
    new_part: str,
    use_default: bool,
) -> bytes:
    root = ET.fromstring(payload)
    override_tag = f"{{{_CONTENT_TYPES_NS}}}Override"
    default_tag = f"{{{_CONTENT_TYPES_NS}}}Default"
    selected = None
    for child in list(root):
        if child.tag == override_tag and child.attrib.get("PartName") == "/" + old_part:
            selected = child
            break
    assert selected is not None, old_part
    content_type = selected.attrib["ContentType"]
    if use_default:
        root.remove(selected)
        extension = PurePosixPath(new_part).suffix.lstrip(".")
        assert extension
        ET.SubElement(
            root,
            default_tag,
            {"Extension": extension, "ContentType": content_type},
        )
    else:
        selected.set("PartName", "/" + new_part)
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _semantic_ooxml_fixture(
    tmp_path: Path,
    kind: str,
    variant: str,
    *,
    hostile: bool,
) -> Path:
    source = _make_valid_office(tmp_path, kind)
    details = _OFFICE_PARTS[kind]
    old_part = str(details["part"])
    old_path = PurePosixPath(old_part)
    if variant == "dat":
        new_part = str(old_path.with_suffix(".dat"))
        use_default = False
    elif variant == "extensionless":
        new_part = str(old_path.with_suffix("")) + "_part"
        use_default = False
    elif variant == "mixed_case":
        new_part = str(old_path.with_suffix(".XmL"))
        use_default = False
    elif variant == "default":
        new_part = str(old_path.with_suffix(".blob"))
        use_default = True
    else:  # pragma: no cover - helper contract
        raise ValueError(variant)

    destination = tmp_path / f"semantic-{kind}-{variant}.{'hostile' if hostile else 'clean'}.{kind}"
    with zipfile.ZipFile(source, "r") as incoming:
        records = [(info, incoming.read(info)) for info in incoming.infolist()]
    with zipfile.ZipFile(destination, "w", compression=zipfile.ZIP_DEFLATED) as outgoing:
        for info, data in records:
            name = info.filename
            if name == old_part:
                name = new_part
                if hostile:
                    data = _insert_dtd(data, str(details["root"]))
            elif name == "[Content_Types].xml":
                data = _content_types_for_renamed_part(
                    data,
                    old_part=old_part,
                    new_part=new_part,
                    use_default=use_default,
                )
            else:
                for relationship_name, old_target in details["relationships"]:
                    if name == relationship_name:
                        new_target = old_target.rsplit("/", 1)[0] + "/" + PurePosixPath(new_part).name if "/" in old_target else PurePosixPath(new_part).name
                        data = data.decode("utf-8").replace(old_target, new_target).encode("utf-8")
            outgoing.writestr(name, data)
    return destination


def _office_loader(kind: str):
    if kind == "docx":
        from leandesk.writer import WriterFrame

        return WriterFrame._load_docx, "docx.Document"
    if kind == "xlsx":
        from leandesk.sheets import SheetsFrame

        return SheetsFrame._load_xlsx, "openpyxl.load_workbook"
    from leandesk.slides import SlidesFrame

    return SlidesFrame._load_pptx, "pptx.Presentation"


def _raw_cp932(hex_bytes: str, *, prefix: bytes = b"", suffix: bytes = b"") -> bytes:
    return b"{\\rtf1\\ansi\\ansicpg932 " + prefix + bytes.fromhex(hex_bytes) + suffix + b"}"


class _MutatingReadHandle:
    """Delegate a file handle while injecting one deterministic source mutation."""

    def __init__(self, handle, source: Path, mode: str) -> None:
        self._handle = handle
        self._source = source
        self._mode = mode
        self._read_injected = False
        self._seek_injected = False

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc, tb):
        return self._handle.__exit__(exc_type, exc, tb)

    def __getattr__(self, name):
        return getattr(self._handle, name)

    def read(self, size: int = -1):
        block = self._handle.read(size)
        if self._mode == "truncate" and not self._read_injected and block:
            self._read_injected = True
            with self._source.open("r+b") as attacker:
                attacker.truncate(max(1, self._source.stat().st_size // 2))
                attacker.flush()
                os.fsync(attacker.fileno())
        return block

    def seek(self, offset: int, whence: int = os.SEEK_SET):
        if self._mode == "same_size_change" and not self._seek_injected and offset == 0:
            self._seek_injected = True
            size = self._source.stat().st_size
            with self._source.open("r+b") as attacker:
                attacker.seek(0)
                attacker.write(b"X" * size)
                attacker.truncate(size)
                attacker.flush()
                os.fsync(attacker.fileno())
        return self._handle.seek(offset, whence)


# ---------------------------------------------------------------------------
# LD-QA-C4-001: bind restore verification/extraction to identical bytes
# ---------------------------------------------------------------------------


def test_restore_valid_a_to_valid_b_path_swap_restores_private_a(tmp_path: Path) -> None:
    archive_a, live, sha_a = _restore_fixture(tmp_path)
    profile_b = _profile(tmp_path / "profile_b", "BACKUP_B_SWAPPED_IN")
    archive_b = _backup(tmp_path / "replacement.ldbackup", profile_b)
    real_verify = backup_service.verify_backup_artifact
    swapped = False

    def verify_private(path, *, require_manifest=False):
        nonlocal swapped
        result = real_verify(path, require_manifest=require_manifest)
        if not swapped:
            os.replace(archive_b, archive_a)
            swapped = True
        return result

    with patch("leandesk.backup_service.verify_backup_artifact", side_effect=verify_private):
        result = restore_backup(archive_a, data_root=live)
    assert result["committed"] is True
    assert result["source_sha256"] == sha_a
    _assert_live(live, "BACKUP_A_VERIFIED")
    assert sha256_file(archive_a) != sha_a


def test_restore_valid_a_to_invalid_path_swap_restores_private_a(tmp_path: Path) -> None:
    archive_a, live, sha_a = _restore_fixture(tmp_path)
    invalid = tmp_path / "invalid.ldbackup"
    invalid.write_bytes(b"NOT_A_BACKUP")
    real_verify = backup_service.verify_backup_artifact

    def verify_private(path, *, require_manifest=False):
        result = real_verify(path, require_manifest=require_manifest)
        os.replace(invalid, archive_a)
        return result

    with patch("leandesk.backup_service.verify_backup_artifact", side_effect=verify_private):
        result = restore_backup(archive_a, data_root=live)
    assert result["source_sha256"] == sha_a
    _assert_live(live, "BACKUP_A_VERIFIED")
    assert archive_a.read_bytes() == b"NOT_A_BACKUP"


def test_restore_source_rename_after_private_copy_does_not_switch_bytes(tmp_path: Path) -> None:
    archive, live, sha = _restore_fixture(tmp_path)
    renamed = tmp_path / "renamed.ldbackup"
    real_verify = backup_service.verify_backup_artifact

    def verify_private(path, *, require_manifest=False):
        result = real_verify(path, require_manifest=require_manifest)
        archive.rename(renamed)
        return result

    with patch("leandesk.backup_service.verify_backup_artifact", side_effect=verify_private):
        result = restore_backup(archive, data_root=live)
    assert result["source_sha256"] == sha
    _assert_live(live, "BACKUP_A_VERIFIED")
    assert renamed.is_file() and not archive.exists()


def test_restore_source_symlink_after_private_copy_does_not_switch_bytes(tmp_path: Path) -> None:
    archive_a, live, sha_a = _restore_fixture(tmp_path)
    profile_b = _profile(tmp_path / "profile_b", "BACKUP_B")
    archive_b = _backup(tmp_path / "backup_b.ldbackup", profile_b)
    displaced = tmp_path / "backup_a.displaced"
    real_verify = backup_service.verify_backup_artifact

    def verify_private(path, *, require_manifest=False):
        result = real_verify(path, require_manifest=require_manifest)
        archive_a.rename(displaced)
        try:
            archive_a.symlink_to(archive_b)
        except (OSError, NotImplementedError):
            displaced.rename(archive_a)
            pytest.skip("symlink creation unavailable")
        return result

    try:
        with patch("leandesk.backup_service.verify_backup_artifact", side_effect=verify_private):
            result = restore_backup(archive_a, data_root=live)
        assert result["source_sha256"] == sha_a
        _assert_live(live, "BACKUP_A_VERIFIED")
        assert archive_a.is_symlink()
    finally:
        if archive_a.is_symlink():
            archive_a.unlink()
        if displaced.exists() and not archive_a.exists():
            displaced.rename(archive_a)


def test_restore_initial_symlink_source_is_refused(tmp_path: Path) -> None:
    archive, live, _ = _restore_fixture(tmp_path)
    linked = tmp_path / "linked.ldbackup"
    try:
        linked.symlink_to(archive)
    except (OSError, NotImplementedError):
        pytest.skip("symlink creation unavailable")
    with pytest.raises(BackupRestoreStateError, match="regular|linked|safely"):
        restore_backup(linked, data_root=live)
    _assert_live(live, "LIVE_OLD")


@pytest.mark.parametrize("mode", ["truncate", "same_size_change"])
def test_restore_concurrent_source_mutation_during_private_copy_is_refused(
    tmp_path: Path, mode: str
) -> None:
    archive, live, _ = _restore_fixture(tmp_path)
    real_open = backup_service._open_bound_regular

    def open_with_mutation(path, *, label):
        handle, before = real_open(path, label=label)
        if label == "backup restore source":
            return _MutatingReadHandle(handle, Path(path), mode), before
        return handle, before

    with patch("leandesk.backup_service._open_bound_regular", side_effect=open_with_mutation):
        with pytest.raises(BackupRestoreStateError, match="changed|copied|source"):
            restore_backup(archive, data_root=live)
    _assert_live(live, "LIVE_OLD")


def test_restore_reported_source_sha_is_exact_private_restored_archive(tmp_path: Path) -> None:
    archive, live, expected_sha = _restore_fixture(tmp_path)
    result = restore_backup(archive, data_root=live)
    assert result["committed"] is True
    assert result["source_sha256"] == expected_sha == sha256_file(archive)
    _assert_live(live, "BACKUP_A_VERIFIED")


# ---------------------------------------------------------------------------
# LD-QA-C4-002: close final verified-temporary/commit gap
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("replacement", ["invalid", "same_bytes_new_inode"])
def test_backup_swap_after_verified_candidate_hash_restores_old_destination(
    tmp_path: Path, replacement: str
) -> None:
    destination, source, old_sha = _old_destination_and_new_profile(tmp_path)
    real_hash = backup_service._hash_regular_path_bound
    swapped = False

    def hash_then_swap(path, *, label):
        nonlocal swapped
        digest, identity = real_hash(path, label=label)
        if label == "verified backup candidate" and not swapped:
            attacker = Path(path).with_name(Path(path).name + ".attacker")
            if replacement == "invalid":
                attacker.write_bytes(b"NOT_A_VALID_BACKUP")
            else:
                shutil.copy2(path, attacker)
            os.replace(attacker, path)
            swapped = True
        return digest, identity

    with patch("leandesk.backup_service._hash_regular_path_bound", side_effect=hash_then_swap):
        with pytest.raises(BackupIntegrityError, match="changed|match"):
            create_backup(destination, data_root=source)
    assert swapped
    _assert_old_destination(destination, old_sha)


def test_backup_posthash_swap_without_previous_destination_leaves_no_target(tmp_path: Path) -> None:
    source = _profile(tmp_path / "profile", "NEW_PROFILE")
    destination = tmp_path / "new.ldbackup"
    real_hash = backup_service._hash_regular_path_bound

    def hash_then_swap(path, *, label):
        digest, identity = real_hash(path, label=label)
        if label == "verified backup candidate":
            attacker = Path(path).with_name(Path(path).name + ".attacker")
            attacker.write_bytes(b"INVALID")
            os.replace(attacker, path)
        return digest, identity

    with patch("leandesk.backup_service._hash_regular_path_bound", side_effect=hash_then_swap):
        with pytest.raises(BackupIntegrityError):
            create_backup(destination, data_root=source)
    assert not destination.exists()


def test_backup_target_inplace_mutation_immediately_after_replace_restores_old(tmp_path: Path) -> None:
    destination, source, old_sha = _old_destination_and_new_profile(tmp_path)
    real_replace = backup_service.os.replace
    injected = False

    def replace_then_mutate(src, dst):
        nonlocal injected
        result = real_replace(src, dst)
        if (
            Path(dst) == destination
            and Path(src).name.startswith(f".{destination.name}.")
            and ".previous-" not in Path(src).name
        ):
            Path(dst).write_bytes(b"POST_REPLACE_MUTATION")
            injected = True
        return result

    with patch("leandesk.backup_service.os.replace", side_effect=replace_then_mutate):
        with pytest.raises(BackupIntegrityError, match="match|changed"):
            create_backup(destination, data_root=source)
    assert injected
    _assert_old_destination(destination, old_sha)


def test_backup_postcommit_manifest_verification_failure_restores_old(tmp_path: Path) -> None:
    destination, source, old_sha = _old_destination_and_new_profile(tmp_path)
    with patch(
        "leandesk.backup_service._backup_integrity.verify_backup_artifact",
        side_effect=BackupIntegrityError("postcommit verification failure"),
    ):
        with pytest.raises(BackupIntegrityError, match="postcommit verification failure"):
            create_backup(destination, data_root=source)
    _assert_old_destination(destination, old_sha)


def test_backup_candidate_install_replace_failure_reactivates_old(tmp_path: Path) -> None:
    destination, source, old_sha = _old_destination_and_new_profile(tmp_path)
    real_replace = backup_service.os.replace

    def fail_candidate_install(src, dst):
        if (
            Path(dst) == destination
            and Path(src).name.startswith(f".{destination.name}.")
            and ".previous-" not in Path(src).name
        ):
            raise OSError("candidate install fault")
        return real_replace(src, dst)

    with patch("leandesk.backup_service.os.replace", side_effect=fail_candidate_install):
        with pytest.raises(BackupIntegrityError, match="backup creation failed"):
            create_backup(destination, data_root=source)
    _assert_old_destination(destination, old_sha)


def test_backup_changed_rollback_copy_is_never_reactivated(tmp_path: Path) -> None:
    destination, source, old_sha = _old_destination_and_new_profile(tmp_path)
    real_hash = backup_service._hash_regular_path_bound
    changed = False

    def corrupt_after_candidate_commit(path, *, label):
        nonlocal changed
        if label == "committed backup destination" and not changed:
            rollback = next(destination.parent.glob(f".{destination.name}.previous-*.tmp"))
            rollback.write_bytes(b"ROLLBACK_TAMPERED")
            Path(path).write_bytes(b"CANDIDATE_TAMPERED")
            changed = True
        return real_hash(path, label=label)

    with patch("leandesk.backup_service._hash_regular_path_bound", side_effect=corrupt_after_candidate_commit):
        with pytest.raises(BackupIntegrityError, match="rollback copy changed|refusing"):
            create_backup(destination, data_root=source)
    assert changed
    # The service refuses to replace the target with attacker-controlled rollback bytes.
    assert sha256_file(destination) != old_sha
    assert any(destination.parent.glob(f".{destination.name}.previous-*.tmp"))


def test_backup_success_reports_exact_committed_destination_hash(tmp_path: Path) -> None:
    destination, source, old_sha = _old_destination_and_new_profile(tmp_path)
    result = create_backup(destination, data_root=source)
    actual = sha256_file(destination)
    assert result["committed"] is True
    assert result["sha256"] == actual
    assert actual != old_sha
    assert verify_backup_artifact(destination, require_manifest=True)["valid"] is True


# ---------------------------------------------------------------------------
# LD-QA-C4-003: semantic OOXML XML classification
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("kind", ["docx", "xlsx", "pptx"])
@pytest.mark.parametrize("variant", ["dat", "extensionless", "mixed_case", "default"])
def test_semantic_xml_dtd_is_rejected_before_downstream_parser(
    tmp_path: Path, kind: str, variant: str
) -> None:
    hostile = _semantic_ooxml_fixture(tmp_path, kind, variant, hostile=True)
    loader, patch_target = _office_loader(kind)
    with patch(patch_target) as downstream:
        with pytest.raises(OOXMLPreflightError, match="forbidden DTD/entity"):
            loader(hostile)
        downstream.assert_not_called()


@pytest.mark.parametrize("kind", ["docx", "xlsx", "pptx"])
@pytest.mark.parametrize("variant", ["dat", "extensionless", "mixed_case", "default"])
def test_clean_semantic_xml_nonstandard_filename_takes_protected_parse_path(
    tmp_path: Path, kind: str, variant: str
) -> None:
    clean = _semantic_ooxml_fixture(tmp_path, kind, variant, hostile=False)
    prepared = prepare_ooxml(clean, kind)
    assert prepared.payload == clean.read_bytes()
    assert prepared.report.xml_bytes > 0


@pytest.mark.parametrize("kind", ["docx", "xlsx", "pptx"])
def test_relationship_content_type_with_non_relationship_filename_is_rejected(
    tmp_path: Path, kind: str
) -> None:
    source = _make_valid_office(tmp_path, kind)
    relation = "_rels/.rels"
    renamed = "_rels/root.dat"
    destination = tmp_path / f"relationship-renamed.{kind}"
    with zipfile.ZipFile(source, "r") as incoming:
        records = [(info, incoming.read(info)) for info in incoming.infolist()]
    with zipfile.ZipFile(destination, "w", zipfile.ZIP_DEFLATED) as outgoing:
        for info, data in records:
            name = info.filename
            if name == relation:
                name = renamed
            elif name == "[Content_Types].xml":
                # Root relationships normally use the Default rels content type.
                root = ET.fromstring(data)
                ET.SubElement(
                    root,
                    f"{{{_CONTENT_TYPES_NS}}}Override",
                    {
                        "PartName": "/" + renamed,
                        "ContentType": "application/vnd.openxmlformats-package.relationships+xml",
                    },
                )
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            outgoing.writestr(name, data)
    with pytest.raises(OOXMLPreflightError, match="missing required parts|Malformed relationship"):
        prepare_ooxml(destination, kind)


# ---------------------------------------------------------------------------
# LD-QA-C4-004: raw CP932 syntax-valued trail bytes
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    ("hex_bytes", "expected"),
    [("955C", "表"), ("815C", "―"), ("817B", "＋"), ("817D", "±")],
)
def test_raw_cp932_syntax_valued_trail_bytes_decode(hex_bytes: str, expected: str) -> None:
    assert rtf_to_plain(_raw_cp932(hex_bytes)) == expected


def test_raw_cp932_trails_do_not_swallow_adjacent_real_controls_or_groups() -> None:
    payload = (
        b"{\\rtf1\\ansi\\ansicpg932 "
        + bytes.fromhex("955C")
        + b"\\b B\\b0 "
        + bytes.fromhex("815C")
        + b"{\\i G\\i0}"
        + bytes.fromhex("817B")
        + b"\\par "
        + bytes.fromhex("817D")
        + b"}"
    )
    assert rtf_to_plain(payload) == "表B―G＋\n±"


def test_escape_origin_lead_preserves_true_hex_control_boundary() -> None:
    assert rtf_to_plain(rb"{\rtf1\ansi\ansicpg936 \'d6\'d0}") == "中"
    assert rtf_to_plain(rb"{\rtf1\ansi\ansicpg932 \'82\b\'a0}") == "あ"


def test_invalid_raw_cp932_lead_before_group_close_fails_controlled() -> None:
    with pytest.raises(RTFFormatError, match="incomplete or invalid"):
        rtf_to_plain(b"{\\rtf1\\ansi\\ansicpg932 " + bytes.fromhex("82") + b"}")


def test_rtf_group_depth_limit_remains_enforced_with_dbcs_logic() -> None:
    payload = b"{\\rtf1\\ansi\\ansicpg932 " + (b"{" * 600) + b"x" + (b"}" * 600) + b"}"
    with pytest.raises(RTFFormatError, match="nesting exceeds"):
        rtf_to_plain(payload)


def test_cp932_syntax_trail_core_save_and_reopen(tmp_path: Path) -> None:
    text = "表―＋±"
    source = tmp_path / "source.rtf"
    source.write_bytes(
        b"{\\rtf1\\ansi\\ansicpg932 "
        + b"".join(bytes.fromhex(value) for value in ("955C", "815C", "817B", "817D"))
        + b"}"
    )
    loaded = read_text_document(source)
    assert loaded.text == text
    destination = tmp_path / "saved.rtf"
    write_text_document(loaded, destination)
    assert destination.read_bytes().isascii()
    assert read_text_document(destination).text == text


@pytest.mark.skipif(not os.environ.get("DISPLAY") and os.name != "nt", reason="GUI test requires Xvfb/display")
def test_writer_gui_cp932_syntax_trail_open_save_as_reopen(tmp_path: Path) -> None:
    import tkinter as tk

    from leandesk.core import AppSettings, RecentFiles, RecoveryStore
    from leandesk.writer import WriterFrame

    expected = "表―＋±"
    source = tmp_path / "source.rtf"
    source.write_bytes(
        b"{\\rtf1\\ansi\\ansicpg932 "
        + b"".join(bytes.fromhex(value) for value in ("955C", "815C", "817B", "817D"))
        + b"}"
    )
    destination = tmp_path / "saved.rtf"
    root = tk.Tk()
    root.withdraw()
    try:
        frame = WriterFrame(root, recent=RecentFiles(tmp_path / "recent.json"), settings=AppSettings())
        frame.recovery = RecoveryStore(tmp_path / "Recovery")
        with patch("leandesk.writer.messagebox.showerror") as showerror:
            assert frame.open_document(source) is True
            assert frame.text.get("1.0", "end-1c") == expected
            with patch("leandesk.writer.filedialog.asksaveasfilename", return_value=str(destination)):
                assert frame.save_as() is True
            showerror.assert_not_called()
        assert read_text_document(destination).text == expected
    finally:
        root.destroy()


@pytest.mark.parametrize(
    ("hex_bytes", "expected"),
    [("955C", "表"), ("815C", "―"), ("817B", "＋"), ("817D", "±")],
)
def test_raw_cp932_syntax_trail_is_readable_by_libreoffice(
    tmp_path: Path, hex_bytes: str, expected: str
) -> None:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        pytest.skip("LibreOffice is not installed")
    source = tmp_path / f"cp932-{hex_bytes}.rtf"
    source.write_bytes(_raw_cp932(hex_bytes))
    outdir = tmp_path / "out"
    profile = tmp_path / "lo-profile"
    outdir.mkdir()
    profile.mkdir()
    result = subprocess.run(
        [
            soffice,
            "--headless",
            "--nologo",
            "--nodefault",
            "--nolockcheck",
            "--norestore",
            f"-env:UserInstallation={profile.resolve().as_uri()}",
            "--convert-to",
            "txt:Text",
            "--outdir",
            str(outdir),
            str(source),
        ],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
    )
    converted = outdir / f"cp932-{hex_bytes}.txt"
    assert result.returncode == 0, (result.stdout + result.stderr).decode("utf-8", errors="replace")
    assert converted.is_file()
    assert expected in converted.read_text(encoding="utf-8-sig", errors="strict")


# ---------------------------------------------------------------------------
# LD-QA-C4-005: raw parent-chain containment for backup sources
# ---------------------------------------------------------------------------


def test_backup_source_linked_parent_is_refused_without_external_capture(tmp_path: Path) -> None:
    external_parent = tmp_path / "external"
    external = _profile(external_parent / "profile", "EXTERNAL_SECRET")
    linked_parent = tmp_path / "linked_parent"
    try:
        linked_parent.symlink_to(external_parent, target_is_directory=True)
    except (OSError, NotImplementedError):
        pytest.skip("symlink creation unavailable")
    destination = tmp_path / "backup.ldbackup"
    with pytest.raises(BackupIntegrityError, match="linked|reparse"):
        create_backup(destination, data_root=linked_parent / external.name)
    assert not destination.exists()


def test_backup_source_linked_root_is_refused(tmp_path: Path) -> None:
    external = _profile(tmp_path / "external", "EXTERNAL_SECRET")
    linked = tmp_path / "linked_profile"
    try:
        linked.symlink_to(external, target_is_directory=True)
    except (OSError, NotImplementedError):
        pytest.skip("symlink creation unavailable")
    with pytest.raises(BackupIntegrityError, match="linked|reparse"):
        create_backup(tmp_path / "backup.ldbackup", data_root=linked)


def test_backup_source_mount_redirection_is_refused(tmp_path: Path) -> None:
    parent = tmp_path / "mounted"
    source = _profile(parent / "profile", "SECRET")
    real_ismount = os.path.ismount

    def simulated_mount(path) -> bool:
        return Path(path) == parent or real_ismount(path)

    with patch("leandesk.backup_service.os.path.ismount", side_effect=simulated_mount):
        with pytest.raises(BackupIntegrityError, match="mount redirection"):
            create_backup(tmp_path / "backup.ldbackup", data_root=source)


@pytest.mark.parametrize("boundary", ["parent", "root"])
def test_backup_source_simulated_windows_reparse_boundary_is_refused(
    tmp_path: Path, boundary: str
) -> None:
    parent = tmp_path / "parent"
    source = _profile(parent / "profile", "SECRET")
    marked = parent if boundary == "parent" else source
    marked_inode = marked.lstat().st_ino
    real = backup_service._info_is_reparse

    def simulated_reparse(info) -> bool:
        return info.st_ino == marked_inode or real(info)

    with patch("leandesk.backup_service._info_is_reparse", side_effect=simulated_reparse):
        with pytest.raises(BackupIntegrityError, match="reparse"):
            create_backup(tmp_path / "backup.ldbackup", data_root=source)


def test_backup_source_parent_replacement_during_manifest_is_detected(tmp_path: Path) -> None:
    parent = tmp_path / "parent"
    source = _profile(parent / "profile", "ORIGINAL")
    displaced = tmp_path / "displaced_parent"
    real_hash = backup_service._hash_profile_file
    replaced = False

    def replace_parent(path, logical, guard):
        nonlocal replaced
        if not replaced:
            parent.rename(displaced)
            _profile(parent / "profile", "ATTACKER")
            replaced = True
        return real_hash(path, logical, guard)

    with patch("leandesk.backup_service._hash_profile_file", side_effect=replace_parent):
        with pytest.raises(BackupIntegrityError, match="containment changed"):
            create_backup(tmp_path / "backup.ldbackup", data_root=source)
    assert replaced
    assert not (tmp_path / "backup.ldbackup").exists()


def test_backup_source_root_replacement_during_member_write_is_detected(tmp_path: Path) -> None:
    source = _profile(tmp_path / "profile", "ORIGINAL")
    displaced = tmp_path / "displaced_profile"
    real_write = backup_service._write_profile_member
    replaced = False

    def replace_root(archive, path, logical, expected):
        nonlocal replaced
        if not replaced:
            source.rename(displaced)
            _profile(source, "ATTACKER")
            replaced = True
        return real_write(archive, path, logical, expected)

    with patch("leandesk.backup_service._write_profile_member", side_effect=replace_root):
        with pytest.raises(BackupIntegrityError, match="containment changed"):
            create_backup(tmp_path / "backup.ldbackup", data_root=source)
    assert replaced
    assert not (tmp_path / "backup.ldbackup").exists()


def test_backup_source_normal_unlinked_tree_still_succeeds(tmp_path: Path) -> None:
    source = _profile(tmp_path / "profile", "NORMAL")
    destination = tmp_path / "backup.ldbackup"
    result = create_backup(destination, data_root=source)
    assert result["committed"] is True
    assert result["sha256"] == sha256_file(destination)


# ---------------------------------------------------------------------------
# Canonical gate and exact source-stage cleanliness
# ---------------------------------------------------------------------------


def test_correction_5_suite_is_mandatory_and_builder_uses_correction_5_environment() -> None:
    runner = (ROOT / "tools" / "run_authoritative_tests.py").read_text(encoding="utf-8")
    assert '"tests/test_correction_5_qa.py"' in runner
    match = re.search(r"MIN_EXPECTED_TESTS\s*=\s*(\d+)", runner)
    assert match and int(match.group(1)) > 280
    builder = (ROOT / "BUILD_LEANDESK_SUITE.ps1").read_text(encoding="utf-8-sig")
    assert '"tests\\test_correction_5_qa.py"' in builder
    assert 'Join-Path $env:TEMP "LeanDesk_0.8.0_Correction_5_Build_Venv"' in builder


def test_correction_5_source_stage_remains_cache_bytecode_and_link_free() -> None:
    from tools.package_cleanliness import scan_tree

    report = scan_tree(ROOT)
    assert report["clean"], report["issues"]

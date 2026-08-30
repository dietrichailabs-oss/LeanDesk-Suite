from __future__ import annotations

import hashlib
import inspect
import io
import json
import os
from pathlib import Path
import re
import shutil
import stat
import subprocess
from unittest.mock import patch
import zipfile

import pytest

from leandesk.backup_integrity import BackupIntegrityError, sha256_file
from leandesk.backup_service import (
    BackupRestoreStateError,
    _prepare_restore_guard,
    _recheck_restore_guard,
    create_backup,
    restore_backup,
)
from leandesk.document_formats import read_text_document
from leandesk.ooxml_preflight import (
    OOXMLLimits,
    OOXMLPreflightCancelled,
    OOXMLPreflightError,
    OOXMLPreflightTimeout,
    prepare_ooxml,
)
from leandesk.rtf_codec import RTFFormatError, rtf_to_plain

ROOT = Path(__file__).resolve().parents[1]
FIXTURES = Path(__file__).resolve().parent / "fixtures" / "correction_4"
RTF_FIXTURES = FIXTURES / "RTF"


def _make_valid_office(tmp_path: Path, kind: str) -> Path:
    path = tmp_path / f"valid.{kind}"
    if kind == "docx":
        from docx import Document

        document = Document()
        document.add_paragraph("VALID_DOCX")
        document.save(path)
    elif kind == "xlsx":
        from openpyxl import Workbook

        workbook = Workbook()
        workbook.active["A1"] = "VALID_XLSX"
        workbook.save(path)
    elif kind == "pptx":
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "VALID_PPTX"
        presentation.save(path)
    else:  # pragma: no cover - test helper contract
        raise ValueError(kind)
    return path


def _rewrite_zip(
    source: Path,
    destination: Path,
    *,
    replacements: dict[str, bytes] | None = None,
    additions: list[tuple[str | zipfile.ZipInfo, bytes]] | None = None,
) -> Path:
    replacements = replacements or {}
    additions = additions or []
    with zipfile.ZipFile(source, "r") as incoming, zipfile.ZipFile(
        destination, "w", compression=zipfile.ZIP_DEFLATED
    ) as outgoing:
        for info in incoming.infolist():
            outgoing.writestr(info, replacements.get(info.filename, incoming.read(info)))
        for name, data in additions:
            outgoing.writestr(name, data)
    return destination


def _without_xml_declaration(payload: bytes) -> bytes:
    return re.sub(br"^\s*<\?xml.*?\?>", b"", payload, count=1, flags=re.DOTALL).lstrip()


def _dtd_payload(original: bytes, mode: str) -> bytes:
    body = _without_xml_declaration(original)
    declaration = '<!DOCTYPE Types [<!ENTITY leandesk "blocked">]>'
    if mode == "after_4k":
        return b'<?xml version="1.0" encoding="UTF-8"?>' + (b" " * 5000) + declaration.encode() + body
    if mode == "utf16":
        text = '<?xml version="1.0" encoding="UTF-16"?>' + declaration + body.decode("utf-8")
        return text.encode("utf-16")
    if mode == "internal_entity":
        return b'<?xml version="1.0" encoding="UTF-8"?>' + declaration.encode() + body
    raise ValueError(mode)


def _dtd_office_package(tmp_path: Path, kind: str, mode: str = "after_4k") -> Path:
    valid = _make_valid_office(tmp_path, kind)
    with zipfile.ZipFile(valid, "r") as archive:
        content_types = archive.read("[Content_Types].xml")
    return _rewrite_zip(
        valid,
        tmp_path / f"dtd-{mode}.{kind}",
        replacements={"[Content_Types].xml": _dtd_payload(content_types, mode)},
    )


def _old_destination_and_new_profile(tmp_path: Path) -> tuple[Path, Path, str]:
    old_profile = tmp_path / "old_profile"
    old_profile.mkdir()
    (old_profile / "state.txt").write_text("KNOWN_GOOD_OLD_BACKUP", encoding="utf-8")
    destination = tmp_path / "selected.ldbackup"
    create_backup(destination, data_root=old_profile)
    before = sha256_file(destination)

    new_profile = tmp_path / "new_profile"
    new_profile.mkdir()
    (new_profile / "state.txt").write_text("NEW_PROFILE_CONTENT", encoding="utf-8")
    return destination, new_profile, before


def _assert_old_destination_preserved(destination: Path, before: str) -> None:
    assert destination.is_file()
    assert sha256_file(destination) == before


# ---------------------------------------------------------------------------
# LD-QA-C3-001: multibyte RTF code-page streams
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    ("codepage", "expected"),
    [(932, "あ"), (936, "中"), (949, "한"), (65001, "Ω")],
)
@pytest.mark.parametrize("form", ["raw", "hex"])
def test_rtf_multibyte_codepage_fixtures_decode_complete_sequences(
    codepage: int, expected: str, form: str
) -> None:
    payload = (RTF_FIXTURES / f"cp{codepage}_{form}.rtf").read_bytes()
    assert rtf_to_plain(payload) == expected


@pytest.mark.parametrize(
    "fixture",
    ["cp932_group_boundary.rtf", "cp932_control_boundary.rtf"],
)
def test_rtf_multibyte_state_survives_group_and_formatting_control_boundaries(fixture: str) -> None:
    assert rtf_to_plain((RTF_FIXTURES / fixture).read_bytes()) == "あ"


def test_rtf_incomplete_multibyte_sequence_fails_controlled() -> None:
    with pytest.raises(RTFFormatError, match="incomplete or invalid"):
        rtf_to_plain(rb"{\rtf1\ansi\ansicpg932 \'82}")


def test_rtf_codepage_change_with_pending_lead_byte_fails_controlled() -> None:
    with pytest.raises(RTFFormatError, match="code-page change"):
        rtf_to_plain(rb"{\rtf1\ansi\ansicpg932 \'82\ansicpg1252 X}")


def test_multicodepage_rtf_core_open_save_and_reopen(tmp_path: Path) -> None:
    source = tmp_path / "multicodepage.rtf"
    source.write_bytes((RTF_FIXTURES / "multicodepage_composite.rtf").read_bytes())
    loaded = read_text_document(source)
    assert loaded.text == "あ中한Ω"

    from leandesk.document_formats import write_text_document

    destination = tmp_path / "saved.rtf"
    write_text_document(loaded, destination)
    assert destination.read_bytes().isascii()
    assert read_text_document(destination).text == "あ中한Ω"


@pytest.mark.skipif(not os.environ.get("DISPLAY") and os.name != "nt", reason="GUI test requires Xvfb/display")
def test_writer_actual_gui_multicodepage_open_save_as_and_reopen(tmp_path: Path) -> None:
    import tkinter as tk

    from leandesk.core import AppSettings, RecentFiles, RecoveryStore
    from leandesk.writer import WriterFrame

    source = tmp_path / "multicodepage.rtf"
    source.write_bytes((RTF_FIXTURES / "multicodepage_composite.rtf").read_bytes())
    destination = tmp_path / "saved.rtf"
    root = tk.Tk()
    root.withdraw()
    try:
        frame = WriterFrame(root, recent=RecentFiles(tmp_path / "recent.json"), settings=AppSettings())
        frame.recovery = RecoveryStore(tmp_path / "Recovery")
        with patch("leandesk.writer.messagebox.showerror") as showerror:
            assert frame.open_document(source) is True
            assert frame.text.get("1.0", "end-1c") == "あ中한Ω"
            with patch("leandesk.writer.filedialog.asksaveasfilename", return_value=str(destination)):
                assert frame.save_as() is True
            showerror.assert_not_called()
        assert read_text_document(destination).text == "あ中한Ω"
    finally:
        root.destroy()


@pytest.mark.parametrize(
    ("codepage", "expected"),
    [(932, "あ"), (936, "中"), (949, "한"), (65001, "Ω")],
)
def test_multibyte_rtf_is_readable_by_libreoffice(tmp_path: Path, codepage: int, expected: str) -> None:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        pytest.skip("LibreOffice is not installed")
    source = tmp_path / f"cp{codepage}.rtf"
    source.write_bytes((RTF_FIXTURES / f"cp{codepage}_hex.rtf").read_bytes())
    outdir = tmp_path / f"out-{codepage}"
    profile = tmp_path / f"lo-{codepage}"
    outdir.mkdir()
    profile.mkdir()
    result = subprocess.run(
        [
            soffice,
            "--headless",
            "--nologo",
            "--nodefault",
            "--nolockcheck",
            "--norestore",
            f"-env:UserInstallation={profile.resolve().as_uri()}",
            "--convert-to",
            "txt:Text",
            "--outdir",
            str(outdir),
            str(source),
        ],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=30,
        check=False,
    )
    converted = outdir / f"cp{codepage}.txt"
    assert result.returncode == 0, (result.stdout + result.stderr).decode("utf-8", errors="replace")
    assert converted.is_file()
    assert expected in converted.read_text(encoding="utf-8-sig", errors="strict")


# ---------------------------------------------------------------------------
# LD-QA-C3-002: parser-level OOXML DTD/entity prohibition
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("kind", ["docx", "xlsx", "pptx"])
@pytest.mark.parametrize("mode", ["after_4k", "utf16", "internal_entity"])
def test_ooxml_parser_level_dtd_rejection_is_offset_and_encoding_independent(
    tmp_path: Path, kind: str, mode: str
) -> None:
    malicious = _dtd_office_package(tmp_path, kind, mode)
    with pytest.raises(OOXMLPreflightError, match="forbidden DTD/entity"):
        prepare_ooxml(malicious, kind)


@pytest.mark.parametrize("kind", ["docx", "xlsx", "pptx"])
def test_ooxml_mixed_case_xml_part_receives_full_xml_security_path(tmp_path: Path, kind: str) -> None:
    valid = _make_valid_office(tmp_path, kind)
    with zipfile.ZipFile(valid, "r") as archive:
        content_types = archive.read("[Content_Types].xml")
    malicious = _rewrite_zip(
        valid,
        tmp_path / f"mixed-case.{kind}",
        additions=[("custom/hostile.XML", _dtd_payload(content_types, "after_4k"))],
    )
    with pytest.raises(OOXMLPreflightError, match="forbidden DTD/entity"):
        prepare_ooxml(malicious, kind)


@pytest.mark.parametrize("kind", ["docx", "xlsx", "pptx"])
def test_ooxml_deep_xml_is_rejected_during_parse(tmp_path: Path, kind: str) -> None:
    valid = _make_valid_office(tmp_path, kind)
    deep = ("<root>" * 80 + "x" + "</root>" * 80).encode("utf-8")
    malicious = _rewrite_zip(valid, tmp_path / f"deep.{kind}", additions=[("custom/deep.xml", deep)])
    with pytest.raises(OOXMLPreflightError, match="nesting exceeds"):
        prepare_ooxml(malicious, kind, limits=OOXMLLimits(max_xml_depth=40))


@pytest.mark.parametrize("kind", ["docx", "xlsx", "pptx"])
def test_ooxml_high_node_xml_is_rejected_during_parse(tmp_path: Path, kind: str) -> None:
    valid = _make_valid_office(tmp_path, kind)
    nodes = ("<root>" + "<n/>" * 1200 + "</root>").encode("utf-8")
    malicious = _rewrite_zip(valid, tmp_path / f"nodes.{kind}", additions=[("custom/nodes.xml", nodes)])
    with pytest.raises(OOXMLPreflightError, match="too many nodes"):
        prepare_ooxml(malicious, kind, limits=OOXMLLimits(max_xml_nodes=800))


def test_ooxml_cancellation_occurs_inside_xml_parser(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")

    def cancel_only_in_parser_callback() -> bool:
        return any(frame.function == "start_element" for frame in inspect.stack())

    with pytest.raises(OOXMLPreflightCancelled):
        prepare_ooxml(valid, "docx", cancel_check=cancel_only_in_parser_callback)


def test_ooxml_deadline_occurs_inside_xml_parser(tmp_path: Path) -> None:
    valid = _make_valid_office(tmp_path, "docx")

    def parser_sensitive_clock() -> float:
        return 10.0 if any(frame.function == "start_element" for frame in inspect.stack()) else 0.0

    with pytest.raises(OOXMLPreflightTimeout):
        prepare_ooxml(valid, "docx", limits=OOXMLLimits(timeout_seconds=1.0), clock=parser_sensitive_clock)


@pytest.mark.parametrize(
    ("kind", "patch_target", "loader"),
    [
        ("docx", "docx.Document", "writer"),
        ("xlsx", "openpyxl.load_workbook", "sheets"),
        ("pptx", "pptx.Presentation", "slides"),
    ],
)
def test_dtd_package_is_rejected_before_third_party_office_parser(
    tmp_path: Path, kind: str, patch_target: str, loader: str
) -> None:
    malicious = _dtd_office_package(tmp_path, kind)
    if loader == "writer":
        from leandesk.writer import WriterFrame

        operation = WriterFrame._load_docx
    elif loader == "sheets":
        from leandesk.sheets import SheetsFrame

        operation = SheetsFrame._load_xlsx
    else:
        from leandesk.slides import SlidesFrame

        operation = SlidesFrame._load_pptx
    with patch(patch_target) as third_party:
        with pytest.raises(OOXMLPreflightError, match="forbidden DTD/entity"):
            operation(malicious)
        third_party.assert_not_called()


# ---------------------------------------------------------------------------
# LD-QA-C3-003: verify temporary backup before destination replacement
# ---------------------------------------------------------------------------


def test_backup_source_mutation_preserves_existing_destination(tmp_path: Path) -> None:
    destination, source, before = _old_destination_and_new_profile(tmp_path)
    from leandesk import backup_service

    real = backup_service._write_profile_member
    mutated = False

    def mutate_before_write(archive, path, logical, expected):
        nonlocal mutated
        if not mutated:
            path.write_text("MUTATED_DURING_BACKUP", encoding="utf-8")
            mutated = True
        return real(archive, path, logical, expected)

    with patch("leandesk.backup_service._write_profile_member", side_effect=mutate_before_write):
        with pytest.raises(BackupIntegrityError):
            create_backup(destination, data_root=source)
    _assert_old_destination_preserved(destination, before)


def test_backup_zip_writer_failure_preserves_existing_destination(tmp_path: Path) -> None:
    destination, source, before = _old_destination_and_new_profile(tmp_path)
    with patch("leandesk.backup_service.zipfile.ZipFile.writestr", side_effect=OSError("zip writer")):
        with pytest.raises(BackupIntegrityError):
            create_backup(destination, data_root=source)
    _assert_old_destination_preserved(destination, before)


@pytest.mark.parametrize("failure", ["manifest mismatch", "semantic validation failure"])
def test_backup_temporary_verification_failure_preserves_existing_destination(
    tmp_path: Path, failure: str
) -> None:
    destination, source, before = _old_destination_and_new_profile(tmp_path)
    with patch(
        "leandesk.backup_service.verify_backup_artifact",
        side_effect=BackupIntegrityError(failure),
    ):
        with pytest.raises(BackupIntegrityError, match=failure):
            create_backup(destination, data_root=source)
    _assert_old_destination_preserved(destination, before)


def test_backup_temporary_flush_failure_preserves_existing_destination(tmp_path: Path) -> None:
    destination, source, before = _old_destination_and_new_profile(tmp_path)
    with patch("leandesk.backup_service.os.fsync", side_effect=OSError("temp flush")):
        with pytest.raises(BackupIntegrityError):
            create_backup(destination, data_root=source)
    _assert_old_destination_preserved(destination, before)


def test_backup_destination_replace_failure_preserves_existing_destination(tmp_path: Path) -> None:
    destination, source, before = _old_destination_and_new_profile(tmp_path)
    with patch("leandesk.backup_service.os.replace", side_effect=OSError("replace")):
        with pytest.raises(BackupIntegrityError):
            create_backup(destination, data_root=source)
    _assert_old_destination_preserved(destination, before)


def test_backup_verified_temporary_swap_before_replace_preserves_existing_destination(tmp_path: Path) -> None:
    destination, source, before = _old_destination_and_new_profile(tmp_path)
    from leandesk import backup_service

    real_recheck = backup_service._recheck_backup_destination

    def replace_verified_temporary(guard) -> None:
        real_recheck(guard)
        candidates = list(destination.parent.glob(f".{destination.name}.*.tmp"))
        assert len(candidates) == 1
        candidate = candidates[0]
        replacement = destination.parent / f".{destination.name}.attacker.tmp"
        replacement.write_bytes(b"not the verified archive")
        os.replace(replacement, candidate)

    with patch(
        "leandesk.backup_service._recheck_backup_destination",
        side_effect=replace_verified_temporary,
    ):
        with pytest.raises(BackupIntegrityError, match="temporary file changed|temporary bytes changed"):
            create_backup(destination, data_root=source)
    _assert_old_destination_preserved(destination, before)


def test_backup_semantic_source_failure_preserves_existing_destination(tmp_path: Path) -> None:
    destination, source, before = _old_destination_and_new_profile(tmp_path)
    (source / "notes.json").write_text("not valid json", encoding="utf-8")
    with pytest.raises(BackupIntegrityError):
        create_backup(destination, data_root=source)
    _assert_old_destination_preserved(destination, before)


def test_backup_verification_occurs_while_old_destination_is_untouched(tmp_path: Path) -> None:
    destination, source, before = _old_destination_and_new_profile(tmp_path)
    from leandesk import backup_service

    real_verify = backup_service.verify_backup_artifact
    observed: list[str] = []

    def verify_temporary(path, *, require_manifest=False):
        observed.append(sha256_file(destination))
        assert Path(path) != destination
        return real_verify(path, require_manifest=require_manifest)

    with patch("leandesk.backup_service.verify_backup_artifact", side_effect=verify_temporary):
        result = create_backup(destination, data_root=source)
    assert observed == [before]
    assert result["committed"] is True
    assert result["sha256"] == sha256_file(destination)
    assert result["sha256"] != before


# ---------------------------------------------------------------------------
# LD-QA-C3-004: raw restore-root containment and identity rechecks
# ---------------------------------------------------------------------------


def _restore_archive(tmp_path: Path) -> Path:
    source = tmp_path / "restore_source"
    source.mkdir()
    (source / "state.txt").write_text("NEW_PROFILE", encoding="utf-8")
    archive = tmp_path / "restore.ldbackup"
    create_backup(archive, data_root=source)
    return archive


def test_restore_rejects_raw_symlink_root_without_touching_external_directory(tmp_path: Path) -> None:
    archive = _restore_archive(tmp_path)
    external = tmp_path / "external"
    external.mkdir()
    sentinel = external / "state.txt"
    sentinel.write_text("EXTERNAL_UNCHANGED", encoding="utf-8")
    linked = tmp_path / "linked_profile"
    try:
        linked.symlink_to(external, target_is_directory=True)
    except (OSError, NotImplementedError):
        pytest.skip("symlink creation unavailable")
    with pytest.raises(BackupRestoreStateError, match="linked|reparse"):
        restore_backup(archive, data_root=linked)
    assert sentinel.read_text(encoding="utf-8") == "EXTERNAL_UNCHANGED"
    assert linked.is_symlink()


def test_restore_rejects_symlink_parent_without_creating_external_profile(tmp_path: Path) -> None:
    archive = _restore_archive(tmp_path)
    external_parent = tmp_path / "external_parent"
    external_parent.mkdir()
    linked_parent = tmp_path / "linked_parent"
    try:
        linked_parent.symlink_to(external_parent, target_is_directory=True)
    except (OSError, NotImplementedError):
        pytest.skip("symlink creation unavailable")
    with pytest.raises(BackupRestoreStateError, match="linked|reparse"):
        restore_backup(archive, data_root=linked_parent / "profile")
    assert not (external_parent / "profile").exists()


def test_restore_rejects_simulated_windows_reparse_target(tmp_path: Path) -> None:
    archive = _restore_archive(tmp_path)
    target = tmp_path / "profile"
    target.mkdir()
    target_inode = target.lstat().st_ino
    from leandesk import backup_service

    real = backup_service._info_is_reparse

    def mark_target(info) -> bool:
        return info.st_ino == target_inode or real(info)

    with patch("leandesk.backup_service._info_is_reparse", side_effect=mark_target):
        with pytest.raises(BackupRestoreStateError, match="reparse"):
            restore_backup(archive, data_root=target)


def test_restore_rejects_mount_redirection_in_parent_chain(tmp_path: Path) -> None:
    archive = _restore_archive(tmp_path)
    parent = tmp_path / "mounted"
    parent.mkdir()
    target = parent / "profile"
    real_ismount = os.path.ismount

    def simulated_mount(path) -> bool:
        return Path(path) == parent or real_ismount(path)

    with patch("leandesk.backup_service.os.path.ismount", side_effect=simulated_mount):
        with pytest.raises(BackupRestoreStateError, match="mount redirection"):
            restore_backup(archive, data_root=target)


def test_restore_parent_identity_recheck_detects_replacement(tmp_path: Path) -> None:
    parent = tmp_path / "parent"
    parent.mkdir()
    target = parent / "profile"
    target.mkdir()
    guard = _prepare_restore_guard(target)
    displaced = tmp_path / "displaced_parent"
    parent.rename(displaced)
    parent.mkdir()
    (parent / "profile").mkdir()
    with pytest.raises(BackupIntegrityError, match="parent changed"):
        _recheck_restore_guard(guard, expected_target=guard.target_identity)


def test_restore_target_replacement_race_is_refused_before_rename(tmp_path: Path) -> None:
    archive = _restore_archive(tmp_path)
    target = tmp_path / "live"
    target.mkdir()
    (target / "state.txt").write_text("OLD_PROFILE", encoding="utf-8")
    external = tmp_path / "external_race"
    external.mkdir()
    (external / "state.txt").write_text("EXTERNAL_UNCHANGED", encoding="utf-8")
    displaced = tmp_path / "displaced_live"

    from leandesk import backup_service

    real_recheck = backup_service._recheck_restore_guard
    calls = 0

    def race_on_second_check(guard, *, expected_target):
        nonlocal calls
        calls += 1
        if calls == 2:
            target.rename(displaced)
            target.symlink_to(external, target_is_directory=True)
        return real_recheck(guard, expected_target=expected_target)

    try:
        with patch("leandesk.backup_service._recheck_restore_guard", side_effect=race_on_second_check):
            with pytest.raises(BackupRestoreStateError):
                restore_backup(archive, data_root=target)
        assert (external / "state.txt").read_text(encoding="utf-8") == "EXTERNAL_UNCHANGED"
        assert displaced.is_dir()
    finally:
        if target.is_symlink():
            target.unlink()
        if displaced.exists() and not target.exists():
            displaced.rename(target)


# ---------------------------------------------------------------------------
# LD-QA-C3-005: exact final staging cleanliness
# ---------------------------------------------------------------------------


def test_package_cleanliness_gate_rejects_cache_and_bytecode(tmp_path: Path) -> None:
    from tools.package_cleanliness import scan_tree

    (tmp_path / ".pytest_cache").mkdir()
    cache = tmp_path / "pkg" / "__pycache__"
    cache.mkdir(parents=True)
    (cache / "module.cpython-313.pyc").write_bytes(b"bytecode")
    report = scan_tree(tmp_path)
    assert report["clean"] is False
    paths = {row["path"] for row in report["issues"]}
    assert ".pytest_cache" in paths
    assert "pkg/__pycache__" in paths


def test_package_cleanliness_gate_accepts_allowlisted_source_shape(tmp_path: Path) -> None:
    from tools.package_cleanliness import scan_tree

    (tmp_path / "leandesk").mkdir()
    (tmp_path / "leandesk" / "module.py").write_text("value = 1\n", encoding="utf-8")
    (tmp_path / "README.md").write_text("clean\n", encoding="utf-8")
    report = scan_tree(tmp_path)
    assert report["clean"] is True
    assert report["issues"] == []


def test_exact_source_stage_is_cache_and_bytecode_free() -> None:
    from tools.package_cleanliness import scan_tree

    report = scan_tree(ROOT)
    assert report["clean"], report["issues"]


def test_correction_4_suite_is_required_and_minimum_gate_is_raised() -> None:
    source = (ROOT / "tools" / "run_authoritative_tests.py").read_text(encoding="utf-8")
    assert '"tests/test_correction_4_qa.py"' in source
    match = re.search(r"MIN_EXPECTED_TESTS\s*=\s*(\d+)", source)
    assert match and int(match.group(1)) >= 250
    assert "PYTHONDONTWRITEBYTECODE" in source
    assert "no:cacheprovider" in source

    builder = (ROOT / "BUILD_LEANDESK_SUITE.ps1").read_text(encoding="utf-8-sig")
    assert '"tests\\test_correction_4_qa.py"' in builder
    assert '"tools\\package_cleanliness.py"' in builder
    assert 'Join-Path $env:TEMP "LeanDesk_0.8.0_Correction_4_Build_Venv"' in builder
    assert 'Join-Path $Root ".build_venv"' not in builder

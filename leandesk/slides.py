from __future__ import annotations

import base64
import io
import json
import uuid
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
import tkinter as tk
from tkinter import filedialog, messagebox, simpledialog, ttk
from typing import Any

from .core import RecentFiles, RecoveryRecord, RecoveryStore, atomic_write_json
from .data_boundary import DataCorruptionError, UnsupportedSchemaVersion, merge_known_and_extra, read_bounded, strict_json_load_bytes
from .compatibility import SLIDES_COMPAT, convert_with_libreoffice
from .ui import COLORS, StatusBar, ResponsiveToolbar, ResponsivePanedwindow
from .save_policy import (
    ImportedSourceProtectionError,
    SavePolicyError,
    UnsupportedSaveFormatError,
    imported_source_for,
    mark_save_boundary,
    validate_destination,
    write_atomically,
)

DECK_FORMAT_VERSION = 1
MAX_IMAGE_BYTES = 20 * 1024 * 1024
MAX_IMAGE_PIXELS = 40_000_000

THEMES = {
    "Midnight": {"background": "#111827", "title": "#f5f0e6", "body": "#c8d1df", "accent": "#5f8dff"},
    "Copper Paper": {"background": "#f7f2e8", "title": "#25222a", "body": "#4d4852", "accent": "#d77d42"},
    "Ocean": {"background": "#073b4c", "title": "#f7fff7", "body": "#d8f3dc", "accent": "#4ecdc4"},
    "Orchid": {"background": "#27163d", "title": "#fff4ff", "body": "#e2cfea", "accent": "#c78af0"},
    "Sunrise": {"background": "#fff6e5", "title": "#31263e", "body": "#5f4b66", "accent": "#f08a5d"},
}

SLIDE_LAYOUTS = ("Title Slide", "Title and Content", "Section Header", "Two Content", "Blank")


@dataclass
class SlideObject:
    object_id: str = field(default_factory=lambda: uuid.uuid4().hex)
    kind: str = "text"
    x: float = 80.0
    y: float = 160.0
    width: float = 360.0
    height: float = 120.0
    text: str = ""
    fill: str = "#4E79A7"
    stroke: str = "#26364A"
    font_size: int = 20
    data: dict[str, Any] = field(default_factory=dict)

    @classmethod
    def from_dict(cls, payload: dict[str, Any]) -> "SlideObject":
        if not isinstance(payload, dict):
            raise DataCorruptionError("Invalid slide object.")
        kind = payload.get("kind", "text")
        if kind not in {"text", "shape", "table", "chart", "image"}:
            raise DataCorruptionError("Unsupported slide object kind.")
        geometry = [payload.get(name, default) for name, default in (("x", 80), ("y", 160), ("width", 360), ("height", 120))]
        if any(isinstance(value, bool) or not isinstance(value, (int, float)) for value in geometry):
            raise DataCorruptionError("Invalid slide object geometry.")
        x, y, width, height = [float(value) for value in geometry]
        if not (0 <= x <= 960 and 0 <= y <= 540 and 1 <= width <= 960 and 1 <= height <= 540):
            raise DataCorruptionError("Slide object geometry is outside the slide.")
        data = payload.get("data", {})
        if not isinstance(data, dict):
            raise DataCorruptionError("Invalid slide object data.")
        return cls(
            object_id=str(payload.get("object_id") or uuid.uuid4().hex), kind=kind,
            x=x, y=y, width=width, height=height, text=str(payload.get("text", ""))[:1_000_000],
            fill=str(payload.get("fill", "#4E79A7")), stroke=str(payload.get("stroke", "#26364A")),
            font_size=max(8, min(96, int(payload.get("font_size", 20)))), data=dict(data),
        )

    def to_dict(self) -> dict[str, Any]:
        return {
            "object_id": self.object_id, "kind": self.kind, "x": self.x, "y": self.y,
            "width": self.width, "height": self.height, "text": self.text,
            "fill": self.fill, "stroke": self.stroke, "font_size": self.font_size,
            "data": dict(self.data),
        }


@dataclass
class SlideModel:
    title: str = "New Slide"
    body: str = "Add your content here."
    theme: str = "Midnight"
    notes: str = ""
    # Kept only for backward source compatibility.  It is never dereferenced during
    # rendering/export and is cleared by the strict loader.
    image_path: str = ""
    image_data: str = ""
    image_media_type: str = ""
    objects: list[SlideObject] = field(default_factory=list)
    layout: str = "Title and Content"
    transition: str = "None"
    extra: dict[str, Any] = field(default_factory=dict, repr=False, compare=False)

    @classmethod
    def from_dict(cls, payload: dict[str, Any]) -> "SlideModel":
        if not isinstance(payload, dict):
            raise DataCorruptionError("Invalid slide data.")
        known = {"title", "body", "theme", "notes", "image_path", "image_data", "image_media_type", "objects", "layout", "transition"}
        title = payload.get("title", "New Slide")
        body = payload.get("body", "Add your content here.")
        theme = payload.get("theme", "Midnight")
        notes = payload.get("notes", "")
        image_data = payload.get("image_data", "")
        image_media_type = payload.get("image_media_type", "image/png" if image_data else "")
        raw_objects = payload.get("objects", [])
        layout = payload.get("layout", "Title and Content")
        transition = payload.get("transition", "None")
        text_fields = (("title", title, 16_384), ("body", body, 4_000_000), ("theme", theme, 256), ("notes", notes, 4_000_000))
        for field_name, value, limit in text_fields:
            if not isinstance(value, str) or len(value) > limit:
                raise DataCorruptionError(f"Invalid slide {field_name}.")
        if not isinstance(image_data, str) or len(image_data) > (MAX_IMAGE_BYTES * 4 // 3 + 16):
            raise DataCorruptionError("Invalid embedded slide image.")
        if not isinstance(image_media_type, str):
            raise DataCorruptionError("Invalid embedded slide image media type.")
        if not isinstance(raw_objects, list) or len(raw_objects) > 2_000 or any(not isinstance(row, dict) for row in raw_objects):
            raise DataCorruptionError("Invalid or oversized slide object list.")
        if layout not in SLIDE_LAYOUTS or transition not in {"None", "Fade", "Push", "Wipe"}:
            raise DataCorruptionError("Invalid slide layout or transition.")
        if image_data:
            if image_media_type != "image/png":
                raise DataCorruptionError("LeanDesk slide images must use normalized PNG data.")
            _validate_embedded_image(image_data)
        elif image_media_type:
            raise DataCorruptionError("Slide image media type is present without image data.")
        return cls(
            title=title,
            body=body,
            theme=theme,
            notes=notes,
            image_path="",  # untrusted external paths are intentionally discarded
            image_data=image_data,
            image_media_type=image_media_type,
            objects=[SlideObject.from_dict(row) for row in raw_objects],
            layout=layout,
            transition=transition,
            extra={k: v for k, v in payload.items() if k not in known},
        )

    def to_dict(self) -> dict[str, Any]:
        return merge_known_and_extra(
            {
                "title": self.title,
                "body": self.body,
                "theme": self.theme,
                "notes": self.notes,
                "image_data": self.image_data,
                "image_media_type": self.image_media_type,
                "objects": [item.to_dict() for item in self.objects],
                "layout": self.layout,
                "transition": self.transition,
            },
            self.extra,
        )


@dataclass
class DeckModel:
    title: str = "Untitled Presentation"
    slides: list[SlideModel] = field(default_factory=lambda: [SlideModel("LeanDesk Slides", "Focused presentations without the clutter.")])
    format_version: int = DECK_FORMAT_VERSION
    extra: dict[str, Any] = field(default_factory=dict, repr=False, compare=False)

    @classmethod
    def from_dict(cls, payload: dict[str, Any]) -> "DeckModel":
        if not isinstance(payload, dict):
            raise DataCorruptionError("Invalid LeanDesk presentation root.")
        version = payload.get("format_version", DECK_FORMAT_VERSION)
        if isinstance(version, bool) or not isinstance(version, int) or version < 1:
            raise DataCorruptionError("Invalid presentation format version.")
        if version > DECK_FORMAT_VERSION:
            raise UnsupportedSchemaVersion(version, DECK_FORMAT_VERSION)
        title = payload.get("title", "Untitled Presentation")
        rows = payload.get("slides", [])
        if not isinstance(title, str) or len(title) > 4096:
            raise DataCorruptionError("Invalid presentation title.")
        if not isinstance(rows, list) or len(rows) > 10_000:
            raise DataCorruptionError("Invalid or oversized slide list.")
        if any(not isinstance(row, dict) for row in rows):
            raise DataCorruptionError("Every presentation slide must be an object.")
        slides = [SlideModel.from_dict(row) for row in rows]
        known = {"title", "slides", "format_version"}
        return cls(
            title,
            slides or [SlideModel()],
            version,
            {k: v for k, v in payload.items() if k not in known},
        )

    def to_dict(self) -> dict[str, Any]:
        return merge_known_and_extra(
            {
                "format_version": DECK_FORMAT_VERSION,
                "title": self.title,
                "slides": [slide.to_dict() for slide in self.slides],
            },
            self.extra,
        )


def _normalize_image_bytes(data: bytes) -> tuple[str, str]:
    if not data or len(data) > MAX_IMAGE_BYTES:
        raise ValueError("Image is empty or exceeds the 20 MiB limit.")
    try:
        from PIL import Image
        with Image.open(io.BytesIO(data)) as probe:
            if probe.width * probe.height > MAX_IMAGE_PIXELS:
                raise ValueError("Image dimensions exceed the safety limit.")
            probe.verify()
        with Image.open(io.BytesIO(data)) as image:
            image.load()
            if image.width * image.height > MAX_IMAGE_PIXELS:
                raise ValueError("Image dimensions exceed the safety limit.")
            if image.mode not in {"RGB", "RGBA"}:
                image = image.convert("RGBA" if "A" in image.getbands() else "RGB")
            output = io.BytesIO()
            image.save(output, format="PNG", optimize=True)
            normalized = output.getvalue()
    except Exception as exc:
        if isinstance(exc, ValueError):
            raise
        raise ValueError("Image could not be validated.") from exc
    if len(normalized) > MAX_IMAGE_BYTES:
        raise ValueError("Normalized image exceeds the 20 MiB limit.")
    return base64.b64encode(normalized).decode("ascii"), "image/png"


def _validate_embedded_image(encoded: str) -> bytes:
    try:
        data = base64.b64decode(encoded, validate=True)
    except Exception as exc:
        raise DataCorruptionError("Embedded slide image is not valid base64.") from exc
    _normalize_image_bytes(data)
    return data


class SlidesFrame(ttk.Frame):
    def __init__(self, master, *, recent: RecentFiles, on_recent_changed=None, on_title_changed=None):
        super().__init__(master)
        self.recent = recent
        self.on_recent_changed = on_recent_changed
        self.on_title_changed = on_title_changed
        self.deck = DeckModel()
        self.current_path: Path | None = None
        self.imported_source_path: Path | None = None
        self.recovery = RecoveryStore()
        self.recovery_id = str(uuid.uuid4())
        self.dirty = False
        self.status_var = tk.StringVar(value="Ready")
        self.position_var = tk.StringVar(value="Slide 1 of 1")
        self.title_var = tk.StringVar()
        self.theme_var = tk.StringVar(value="Midnight")
        self.layout_var = tk.StringVar(value="Title and Content")
        self.transition_var = tk.StringVar(value="None")
        self.selected_object_id: str | None = None
        self._image_ref = None
        self._build_ui()
        self.refresh_slides(0)

    def _build_ui(self) -> None:
        ribbon = ResponsiveToolbar(self, bg=COLORS["panel"], height=92, highlightbackground=COLORS["line"], highlightthickness=1)
        ribbon.pack(fill="x")
        ribbon.pack_propagate(False)
        for label, command in (
            ("New", self.new_deck), ("Open", self.open_deck), ("Save", self.save), ("Save As", self.save_as),
            ("Add Slide", self.add_slide), ("Duplicate", self.duplicate_slide), ("Delete", self.delete_slide),
            ("Move Up", lambda: self.move_slide(-1)), ("Move Down", lambda: self.move_slide(1)),
            ("Text Box", self.add_text_box), ("Shape", self.add_shape), ("Table", self.add_table),
            ("Chart", self.add_chart), ("Add Image", self.add_image), ("Present", self.present), ("Export PPTX", self.export_pptx),
        ):
            ttk.Button(ribbon, text=label, command=command).pack(side="left", padx=3, pady=17)
        tk.Label(ribbon, text="SLIDES", bg=COLORS["panel"], fg=COLORS["amber"], font=("Segoe UI Bold", 14)).pack(side="right", padx=16)

        body = ResponsivePanedwindow(self, orient="horizontal")
        body.pack(fill="both", expand=True)
        left = ttk.Frame(body, style="Panel.TFrame", width=210)
        center = ttk.Frame(body)
        right = ttk.Frame(body, style="Panel.TFrame", width=300)
        body.add(left, weight=1)
        body.add(center, weight=5)
        body.add(right, weight=2)

        tk.Label(left, text="SLIDES", bg=COLORS["panel"], fg=COLORS["amber"], font=("Segoe UI Semibold", 10)).pack(anchor="w", padx=12, pady=(12, 6))
        self.slide_list = tk.Listbox(left, bg=COLORS["field"], fg=COLORS["field_text"], selectbackground=COLORS["selection"], selectforeground=COLORS["button_active_text"], relief="flat", bd=0, font=("Segoe UI", 10), activestyle="none")
        self.slide_list.pack(fill="both", expand=True, padx=10, pady=(0, 10))
        self.slide_list.bind("<<ListboxSelect>>", self.slide_selected)

        self.canvas = tk.Canvas(center, bg=COLORS["workspace"], highlightthickness=0)
        self.canvas.pack(fill="both", expand=True)
        self.canvas.bind("<Configure>", lambda _e: self.render_slide())

        tk.Label(right, text="SLIDE CONTENT", bg=COLORS["panel"], fg=COLORS["amber"], font=("Segoe UI Semibold", 10)).pack(anchor="w", padx=12, pady=(12, 5))
        tk.Label(right, text="Title", bg=COLORS["panel"], fg=COLORS["muted"]).pack(anchor="w", padx=12)
        title_entry = ttk.Entry(right, textvariable=self.title_var)
        title_entry.pack(fill="x", padx=12, pady=(2, 8))
        title_entry.bind("<KeyRelease>", self.form_changed)
        tk.Label(right, text="Body", bg=COLORS["panel"], fg=COLORS["muted"]).pack(anchor="w", padx=12)
        self.body_text = tk.Text(right, height=9, wrap="word", bg=COLORS["field"], fg=COLORS["field_text"], insertbackground=COLORS["field_text"], relief="flat", padx=8, pady=8)
        self.body_text.pack(fill="x", padx=12, pady=(2, 8))
        self.body_text.bind("<KeyRelease>", self.form_changed)
        tk.Label(right, text="Theme", bg=COLORS["panel"], fg=COLORS["muted"]).pack(anchor="w", padx=12)
        theme = ttk.Combobox(right, textvariable=self.theme_var, values=tuple(THEMES), state="readonly")
        theme.pack(fill="x", padx=12, pady=(2, 8))
        theme.bind("<<ComboboxSelected>>", self.form_changed)
        tk.Label(right, text="Layout / transition", bg=COLORS["panel"], fg=COLORS["muted"]).pack(anchor="w", padx=12)
        layout = ttk.Combobox(right, textvariable=self.layout_var, values=SLIDE_LAYOUTS, state="readonly")
        layout.pack(fill="x", padx=12, pady=(2, 4))
        layout.bind("<<ComboboxSelected>>", self.form_changed)
        transition = ttk.Combobox(right, textvariable=self.transition_var, values=("None", "Fade", "Push", "Wipe"), state="readonly")
        transition.pack(fill="x", padx=12, pady=(2, 8))
        transition.bind("<<ComboboxSelected>>", self.form_changed)
        tk.Label(right, text="Speaker notes", bg=COLORS["panel"], fg=COLORS["muted"]).pack(anchor="w", padx=12)
        self.notes_text = tk.Text(right, height=8, wrap="word", bg=COLORS["field"], fg=COLORS["field_text"], insertbackground=COLORS["field_text"], relief="flat", padx=8, pady=8)
        self.notes_text.pack(fill="both", expand=True, padx=12, pady=(2, 12))
        self.notes_text.bind("<KeyRelease>", self.form_changed)

        status = StatusBar(self)
        status.pack(fill="x")
        status.add_left(self.status_var)
        status.add_right(self.position_var, muted=True)

    def current_index(self) -> int:
        selection = self.slide_list.curselection()
        return selection[0] if selection else 0

    def current_slide(self) -> SlideModel:
        return self.deck.slides[self.current_index()]

    def refresh_slides(self, select: int | None = None) -> None:
        self.slide_list.delete(0, "end")
        for index, slide in enumerate(self.deck.slides, 1):
            name = slide.title.strip() or "Untitled Slide"
            self.slide_list.insert("end", f"{index}. {name[:28]}")
        index = min(select if select is not None else self.current_index(), len(self.deck.slides) - 1)
        self.slide_list.selection_set(index)
        self.slide_list.activate(index)
        self.load_form(index)
        self.position_var.set(f"Slide {index + 1} of {len(self.deck.slides)}")
        self._update_title()

    def load_form(self, index: int) -> None:
        slide = self.deck.slides[index]
        self.title_var.set(slide.title)
        self.theme_var.set(slide.theme if slide.theme in THEMES else "Midnight")
        self.layout_var.set(slide.layout)
        self.transition_var.set(slide.transition)
        self.body_text.delete("1.0", "end")
        self.body_text.insert("1.0", slide.body)
        self.notes_text.delete("1.0", "end")
        self.notes_text.insert("1.0", slide.notes)
        self.render_slide()

    def slide_selected(self, _event=None) -> None:
        selection = self.slide_list.curselection()
        if selection:
            self.load_form(selection[0])
            self.position_var.set(f"Slide {selection[0] + 1} of {len(self.deck.slides)}")

    def form_changed(self, _event=None) -> None:
        if not self.deck.slides:
            return
        slide = self.current_slide()
        slide.title = self.title_var.get()
        slide.body = self.body_text.get("1.0", "end-1c")
        slide.notes = self.notes_text.get("1.0", "end-1c")
        slide.theme = self.theme_var.get()
        slide.layout = self.layout_var.get()
        slide.transition = self.transition_var.get()
        self.dirty = True
        index = self.current_index()
        self.slide_list.delete(index)
        self.slide_list.insert(index, f"{index + 1}. {(slide.title.strip() or 'Untitled Slide')[:28]}")
        self.slide_list.selection_set(index)
        self.render_slide()
        self._update_title()
        self._save_recovery()

    def _save_recovery(self) -> None:
        try:
            self.recovery.save(
                RecoveryRecord(
                    self.recovery_id,
                    "Slides",
                    self.deck.title,
                    str(self.current_path or ""),
                    datetime.now().isoformat(timespec="seconds"),
                    self.deck.to_dict(),
                )
            )
        except Exception:
            self.status_var.set("Recovery copy could not be updated")

    def render_slide(self) -> None:
        self.canvas.delete("all")
        if not self.deck.slides:
            return
        slide = self.current_slide()
        theme = THEMES.get(slide.theme, THEMES["Midnight"])
        width = max(1, self.canvas.winfo_width())
        height = max(1, self.canvas.winfo_height())
        ratio = max(0.001, min(max(1, width - 24) / 960, max(1, height - 24) / 540))
        sw, sh = int(960 * ratio), int(540 * ratio)
        x1, y1 = (width - sw) // 2, (height - sh) // 2
        x2, y2 = x1 + sw, y1 + sh
        self.canvas.create_rectangle(x1, y1, x2, y2, fill=theme["background"], outline="#556070", width=1)
        self.canvas.create_rectangle(x1, y1, x1 + max(8, int(12 * ratio)), y2, fill=theme["accent"], outline="")
        self.canvas.create_text(x1 + int(72 * ratio), y1 + int(105 * ratio), text=slide.title or "Untitled Slide", fill=theme["title"], anchor="nw", width=int(780 * ratio), font=("Segoe UI Semibold", max(16, int(31 * ratio))))
        self.canvas.create_text(x1 + int(75 * ratio), y1 + int(210 * ratio), text=slide.body, fill=theme["body"], anchor="nw", width=int(770 * ratio), font=("Segoe UI", max(10, int(18 * ratio))))
        for item in slide.objects:
            left = x1 + int(item.x * ratio)
            top = y1 + int(item.y * ratio)
            right = left + int(item.width * ratio)
            bottom = top + int(item.height * ratio)
            outline = theme["accent"] if item.object_id == self.selected_object_id else item.stroke
            if item.kind == "text":
                self.canvas.create_text(left, top, text=item.text, fill=theme["body"], anchor="nw", width=max(20, right - left), font=("Segoe UI", max(8, int(item.font_size * ratio))))
            elif item.kind == "table":
                rows = max(1, int(item.data.get("rows", 2)))
                cols = max(1, int(item.data.get("cols", 2)))
                values = item.data.get("values", [])
                for row in range(rows):
                    for col in range(cols):
                        cx1 = left + (right - left) * col / cols
                        cy1 = top + (bottom - top) * row / rows
                        cx2 = left + (right - left) * (col + 1) / cols
                        cy2 = top + (bottom - top) * (row + 1) / rows
                        self.canvas.create_rectangle(cx1, cy1, cx2, cy2, fill=theme["background"], outline=outline)
                        value = values[row][col] if isinstance(values, list) and row < len(values) and isinstance(values[row], list) and col < len(values[row]) else ""
                        self.canvas.create_text(cx1 + 5, (cy1 + cy2) / 2, text=value, anchor="w", fill=theme["body"], font=("Segoe UI", max(7, int(12 * ratio))))
            elif item.kind == "chart":
                values = item.data.get("values", [3, 5, 2])
                numeric = [float(value) for value in values if isinstance(value, (int, float))] or [1.0]
                maximum = max(numeric) or 1
                bar_width = (right - left) / max(1, len(numeric))
                for index, value in enumerate(numeric):
                    bar_height = (bottom - top - 20) * value / maximum
                    self.canvas.create_rectangle(left + index * bar_width + 4, bottom - bar_height, left + (index + 1) * bar_width - 4, bottom, fill=item.fill, outline=outline)
            else:
                self.canvas.create_rectangle(left, top, right, bottom, fill=item.fill, outline=outline, width=2)
                if item.text:
                    self.canvas.create_text((left + right) / 2, (top + bottom) / 2, text=item.text, fill="#FFFFFF", font=("Segoe UI", max(8, int(item.font_size * ratio))))
        self.canvas.create_oval(x2 - int(150 * ratio), y2 - int(120 * ratio), x2 - int(50 * ratio), y2 - int(20 * ratio), outline=theme["accent"], width=max(2, int(3 * ratio)))
        if slide.image_data:
            try:
                from PIL import Image, ImageTk
                image = Image.open(io.BytesIO(_validate_embedded_image(slide.image_data)))
                image.thumbnail((int(300 * ratio), int(170 * ratio)))
                self._image_ref = ImageTk.PhotoImage(image)
                self.canvas.create_image(x2 - int(190 * ratio), y2 - int(125 * ratio), image=self._image_ref, anchor="se")
            except Exception:
                self._image_ref = None

    def _update_title(self) -> None:
        name = self.current_path.name if self.current_path else self.deck.title
        if self.on_title_changed:
            self.on_title_changed(f"Slides — {name}", self.dirty)

    def confirm_discard(self) -> bool:
        if not self.dirty:
            return True
        answer = messagebox.askyesnocancel("LeanDesk Slides", "Save presentation changes?", parent=self)
        if answer is None:
            return False
        return self.save() if answer else True

    def new_deck(self) -> bool:
        if not self.confirm_discard():
            return False
        self.deck = DeckModel()
        self.current_path = None
        self.imported_source_path = None
        self.recovery.delete(self.recovery_id)
        self.recovery_id = str(uuid.uuid4())
        self.dirty = False
        self.refresh_slides(0)
        return True

    def open_deck(self, path: Path | None = None) -> bool:
        if not self.confirm_discard():
            return False
        if path is None:
            value = filedialog.askopenfilename(parent=self, filetypes=(("LeanDesk presentation", "*.ldeck"), ("PowerPoint/OpenDocument", "*.pptx *.ppt *.pptm *.pps *.ppsx *.odp *.otp"), ("Legacy/other presentations", "*.sxi *.key"), ("All files", "*.*")))
            if not value:
                return False
            path = Path(value)
        try:
            suffix = path.suffix.lower()
            if suffix == ".pptx":
                self.deck = self._load_pptx(path)
            elif suffix in SLIDES_COMPAT:
                converted = convert_with_libreoffice(path, "Slides")
                self.deck = self._load_pptx(converted.as_file(), title=path.stem)
                self.status_var.set(converted.note)
            else:
                payload = strict_json_load_bytes(read_bounded(path, limit=128 * 1024 * 1024))
                self.deck = DeckModel.from_dict(payload)
        except Exception as exc:
            messagebox.showerror("LeanDesk Slides", f"Could not open presentation.\n\n{exc}", parent=self)
            return False
        self.current_path = path
        self.imported_source_path = imported_source_for("Slides", path)
        self.recovery.delete(self.recovery_id)
        self.recovery_id = str(uuid.uuid4())
        self.dirty = False
        self.refresh_slides(0)
        self.recent.add(path, "Slides")
        if self.on_recent_changed:
            self.on_recent_changed()
        return True

    def _protected_import_source(self) -> Path | None:
        explicit = getattr(self, "imported_source_path", None)
        if explicit is not None:
            return Path(explicit)
        return imported_source_for("Slides", getattr(self, "current_path", None))

    @mark_save_boundary
    def save(self) -> bool:
        if self._protected_import_source() is not None:
            try:
                proceed = messagebox.askyesno(
                    "Original File Protected",
                    "This presentation was imported from another format. LeanDesk will not overwrite the original because unsupported content could be lost.\n\nSave a new copy instead?",
                    parent=self,
                )
            except Exception:
                return False
            return self.save_as() if proceed else False
        return self.save_as() if self.current_path is None else self._write(self.current_path)

    @mark_save_boundary
    def save_as(self) -> bool:
        value = filedialog.asksaveasfilename(parent=self, defaultextension=".ldeck", filetypes=(("LeanDesk presentation", "*.ldeck"), ("PowerPoint", "*.pptx")))
        if not value:
            return False
        return self._write(Path(value))

    @mark_save_boundary
    def _write(self, path: Path) -> bool:
        try:
            destination = validate_destination("Slides", path, imported_source=self._protected_import_source())
            suffix = destination.suffix.lower()

            def produce(temporary: Path) -> None:
                if suffix == ".pptx":
                    self._save_pptx(temporary)
                else:
                    atomic_write_json(temporary, self.deck.to_dict())

            write_atomically(destination, produce)
        except (ImportedSourceProtectionError, UnsupportedSaveFormatError, SavePolicyError) as exc:
            messagebox.showwarning("LeanDesk Slides", str(exc), parent=self)
            return False
        except Exception as exc:
            messagebox.showerror("LeanDesk Slides", f"Could not save presentation.\n\n{exc}", parent=self)
            return False
        self.current_path = destination
        self.imported_source_path = None
        self.dirty = False
        self.recovery.delete(self.recovery_id)
        self.recent.add(destination, "Slides")
        if self.on_recent_changed:
            self.on_recent_changed()
        self._update_title()
        self.status_var.set(f"Saved {destination.name}")
        return True

    def add_slide(self) -> None:
        index = self.current_index() + 1
        theme = self.current_slide().theme if self.deck.slides else "Midnight"
        self.deck.slides.insert(index, SlideModel(theme=theme))
        self.dirty = True
        self.refresh_slides(index)

    def duplicate_slide(self) -> None:
        slide = self.current_slide()
        copy = SlideModel.from_dict(slide.to_dict())
        copy.title = f"{copy.title} Copy"
        index = self.current_index() + 1
        self.deck.slides.insert(index, copy)
        self.dirty = True
        self.refresh_slides(index)

    def delete_slide(self) -> None:
        if len(self.deck.slides) == 1:
            messagebox.showinfo("LeanDesk Slides", "A presentation must keep at least one slide.", parent=self)
            return
        index = self.current_index()
        if messagebox.askyesno("LeanDesk Slides", "Delete the selected slide?", parent=self):
            self.deck.slides.pop(index)
            self.dirty = True
            self.refresh_slides(max(0, index - 1))

    def move_slide(self, delta: int) -> None:
        index = self.current_index()
        target = index + delta
        if 0 <= target < len(self.deck.slides):
            self.deck.slides[index], self.deck.slides[target] = self.deck.slides[target], self.deck.slides[index]
            self.dirty = True
            self.refresh_slides(target)

    def add_image(self) -> None:
        value = filedialog.askopenfilename(parent=self, filetypes=(("Images", "*.png *.jpg *.jpeg *.gif *.bmp *.webp"),))
        if value:
            try:
                encoded, media_type = _normalize_image_bytes(read_bounded(Path(value), limit=MAX_IMAGE_BYTES))
            except Exception as exc:
                messagebox.showerror("LeanDesk Slides", f"Could not add the image.\n\n{exc}", parent=self)
                return
            slide = self.current_slide()
            slide.image_path = ""
            slide.image_data = encoded
            slide.image_media_type = media_type
            self.dirty = True
            self.render_slide()
            self._save_recovery()

    def add_object(self, kind: str, **values: Any) -> SlideObject:
        item = SlideObject(kind=kind, **values)
        self.current_slide().objects.append(item)
        self.selected_object_id = item.object_id
        self.dirty = True
        self.render_slide()
        self._save_recovery()
        return item

    def add_text_box(self) -> None:
        text = simpledialog.askstring("Text Box", "Text:", parent=self)
        if text is not None:
            self.add_object("text", text=text)

    def add_shape(self) -> None:
        self.add_object("shape", text="Shape")

    def add_table(self) -> None:
        self.add_object("table", width=480, height=180, data={"rows": 3, "cols": 3, "values": [["" for _ in range(3)] for _ in range(3)]})

    def add_chart(self) -> None:
        self.add_object("chart", width=480, height=240, data={"categories": ["A", "B", "C"], "values": [3, 5, 2]})

    def move_object(self, object_id: str, dx: float, dy: float) -> None:
        item = next(item for item in self.current_slide().objects if item.object_id == object_id)
        item.x = max(0, min(960 - item.width, item.x + dx))
        item.y = max(0, min(540 - item.height, item.y + dy))
        self.dirty = True
        self.render_slide()

    def resize_object(self, object_id: str, width: float, height: float) -> None:
        item = next(item for item in self.current_slide().objects if item.object_id == object_id)
        item.width = max(20, min(960 - item.x, width))
        item.height = max(20, min(540 - item.y, height))
        self.dirty = True
        self.render_slide()

    def reorder_object(self, object_id: str, delta: int) -> None:
        objects = self.current_slide().objects
        index = next(index for index, item in enumerate(objects) if item.object_id == object_id)
        target = max(0, min(len(objects) - 1, index + delta))
        objects.insert(target, objects.pop(index))
        self.dirty = True
        self.render_slide()

    def present(self) -> None:
        window = tk.Toplevel(self)
        window.title("LeanDesk Slides — Presenter")
        window.configure(bg="#000000")
        window.geometry("1200x720")
        canvas = tk.Canvas(window, bg="#000000", highlightthickness=0)
        canvas.pack(fill="both", expand=True)
        index = [self.current_index()]

        def render():
            canvas.delete("all")
            slide = self.deck.slides[index[0]]
            theme = THEMES.get(slide.theme, THEMES["Midnight"])
            width, height = max(800, canvas.winfo_width()), max(450, canvas.winfo_height())
            margin = 35
            canvas.create_rectangle(margin, margin, width - margin, height - margin, fill=theme["background"], outline="")
            canvas.create_text(margin + 70, margin + 90, text=slide.title, fill=theme["title"], anchor="nw", width=width - 2 * margin - 140, font=("Segoe UI Semibold", 34))
            canvas.create_text(margin + 75, margin + 220, text=slide.body, fill=theme["body"], anchor="nw", width=width - 2 * margin - 150, font=("Segoe UI", 22))
            canvas.create_text(width - margin - 20, height - margin - 15, text=f"{index[0] + 1} / {len(self.deck.slides)}", fill=theme["accent"], anchor="se", font=("Segoe UI", 12))

        def step(delta):
            index[0] = max(0, min(len(self.deck.slides) - 1, index[0] + delta))
            render()

        window.bind("<Right>", lambda _e: step(1))
        window.bind("<space>", lambda _e: step(1))
        window.bind("<Left>", lambda _e: step(-1))
        window.bind("<Escape>", lambda _e: window.destroy())
        canvas.bind("<Configure>", lambda _e: render())
        window.after(50, render)

    def export_pptx(self) -> bool:
        value = filedialog.asksaveasfilename(parent=self, defaultextension=".pptx", filetypes=(("PowerPoint", "*.pptx"),))
        return bool(value) and self._write(Path(value))

    def _save_pptx(self, path: Path) -> None:
        try:
            from pptx import Presentation
            from pptx.chart.data import ChartData
            from pptx.enum.chart import XL_CHART_TYPE
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor
            from pptx.util import Inches, Pt
        except ImportError as exc:
            raise RuntimeError("PPTX support requires python-pptx.") from exc
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        for model in self.deck.slides:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            theme = THEMES.get(model.theme, THEMES["Midnight"])
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = RGBColor.from_string(theme["background"].lstrip("#"))
            title = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.2), Inches(1.2))
            p = title.text_frame.paragraphs[0]
            p.text = model.title
            p.font.size = Pt(30)
            p.font.bold = True
            p.font.color.rgb = RGBColor.from_string(theme["title"].lstrip("#"))
            body = slide.shapes.add_textbox(Inches(1.05), Inches(2.4), Inches(10.9), Inches(3.4))
            p = body.text_frame.paragraphs[0]
            p.text = model.body
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor.from_string(theme["body"].lstrip("#"))
            if model.image_data:
                try:
                    image_bytes = _validate_embedded_image(model.image_data)
                    slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(9.0), Inches(4.7), width=Inches(3.2))
                except Exception:
                    pass
            if model.notes:
                try:
                    slide.notes_slide.notes_text_frame.text = model.notes
                except Exception:
                    pass
            for item in model.objects:
                left, top = Inches(item.x / 72), Inches(item.y / 72)
                width, height = Inches(item.width / 72), Inches(item.height / 72)
                if item.kind == "text":
                    shape = slide.shapes.add_textbox(left, top, width, height)
                    shape.text_frame.text = item.text
                    shape.text_frame.paragraphs[0].font.size = Pt(item.font_size)
                elif item.kind == "shape":
                    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
                    shape.text = item.text
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string(item.fill.lstrip("#"))
                elif item.kind == "table":
                    rows, cols = max(1, int(item.data.get("rows", 2))), max(1, int(item.data.get("cols", 2)))
                    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                    values = item.data.get("values", [])
                    for row in range(rows):
                        for col in range(cols):
                            if isinstance(values, list) and row < len(values) and isinstance(values[row], list) and col < len(values[row]):
                                table.cell(row, col).text = str(values[row][col])
                elif item.kind == "chart":
                    chart_data = ChartData()
                    categories = item.data.get("categories", ["A", "B", "C"])
                    values = item.data.get("values", [3, 5, 2])
                    chart_data.categories = [str(value) for value in categories]
                    chart_data.add_series(item.text or "Series 1", [float(value) for value in values])
                    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data)
        presentation.save(path)

    @staticmethod
    def _load_pptx(path, *, title: str | None = None) -> DeckModel:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as exc:
            raise RuntimeError("PPTX support requires python-pptx.") from exc
        from .ooxml_preflight import prepare_ooxml
        prepared = prepare_ooxml(path, "pptx")
        source = Presentation(prepared.open())
        slides: list[SlideModel] = []
        for source_slide in source.slides:
            texts = [shape.text.strip() for shape in source_slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            slide_title = texts[0] if texts else "Slide"
            body = "\n\n".join(texts[1:])
            image_data = ""
            image_media_type = ""
            for shape in source_slide.shapes:
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_data, image_media_type = _normalize_image_bytes(shape.image.blob)
                        break
                    except Exception:
                        continue
            notes = ""
            try:
                notes = source_slide.notes_slide.notes_text_frame.text or ""
            except Exception:
                pass
            slides.append(SlideModel(slide_title, body, "Midnight", notes, "", image_data, image_media_type))
        source_title = title or (Path(path).stem if isinstance(path, (str, Path)) else "Imported Presentation")
        return DeckModel(source_title, slides or [SlideModel()])

    def recover_record(self, record: RecoveryRecord) -> None:
        if record.module != "Slides" or not self.confirm_discard():
            return
        self.deck = DeckModel.from_dict(record.payload)
        self.current_path = Path(record.original_path) if record.original_path else None
        self.imported_source_path = imported_source_for("Slides", self.current_path)
        self.recovery_id = record.recovery_id
        self.dirty = True
        self.refresh_slides(0)
        self._update_title()
